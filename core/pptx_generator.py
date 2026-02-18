"""
Potomac Brand-Compliant PPTX Generator
=======================================
Builds Potomac-branded presentations using python-pptx.

Brand:
  - Yellow:  #FEC00F  (headers, accents, CTA highlights)
  - Dark:    #212121  (dark slide backgrounds, header bars)
  - Teal:    #00DED1  (secondary accents, horizontal rules)
  - White:   #FFFFFF  (content backgrounds, text on dark)
  - Fonts:   Rajdhani (headers, ALL CAPS) | Quicksand (body text)
  - Layout:  LAYOUT_WIDE — 13.33 x 7.5 inches (widescreen)

Usage:
    gen = PotomacPPTXGenerator()
    path = gen.generate_from_outline(outline_dict, output_path)
"""

import os
import logging
from pathlib import Path
from typing import Optional, List, Dict, Any

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

logger = logging.getLogger(__name__)

# ── Brand Color Constants ─────────────────────────────────────────────────────
YELLOW   = RGBColor(0xFE, 0xC0, 0x0F)
DARK     = RGBColor(0x21, 0x21, 0x21)
TEAL     = RGBColor(0x00, 0xDE, 0xD1)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF5, 0xF5, 0xF5)
MID_GRAY = RGBColor(0x55, 0x55, 0x55)
GREEN    = RGBColor(0x22, 0xC5, 0x5E)
RED      = RGBColor(0xEF, 0x44, 0x44)
DARK_333 = RGBColor(0x33, 0x33, 0x33)
GRAY_999 = RGBColor(0x99, 0x99, 0x99)

# MSO_AUTO_SHAPE_TYPE integers
RECT       = 1   # Rectangle
OVAL       = 9   # Ellipse / circle
ROUND_RECT = 5   # Rounded rectangle

# Slide dimensions
SLIDE_W = 13.33
SLIDE_H = 7.5

# Default asset base — sibling projects folder
_DEFAULT_ASSETS = Path(__file__).parent.parent.parent / "skills" / "skills" / "potomac-pptx"
POTOMAC_ASSETS = Path(os.getenv("POTOMAC_ASSETS_PATH", str(_DEFAULT_ASSETS)))


# ─────────────────────────────────────────────────────────────────────────────
class PotomacPPTXGenerator:
    """Generates Potomac-branded .pptx files from a slide outline dict."""

    def __init__(self, assets_base: Optional[str] = None):
        self.assets = Path(assets_base) if assets_base else POTOMAC_ASSETS
        self.logos  = self.assets / "brand-assets" / "logos"
        self.images = self.assets / "adobe-assets" / "images"

    # ─── low-level drawing helpers ───────────────────────────────────────────

    def _rect(self, slide, x, y, w, h, color: RGBColor, shape_type: int = RECT):
        """Add a solid filled rectangle (no border)."""
        shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        return shape

    def _text(
        self, slide, text: str,
        x, y, w, h,
        size: float = 12,
        color: RGBColor = DARK,
        font: str = "Quicksand",
        bold: bool = False,
        italic: bool = False,
        align: PP_ALIGN = PP_ALIGN.LEFT,
        word_wrap: bool = True,
    ):
        """Add a standalone text box."""
        txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = txBox.text_frame
        tf.word_wrap = word_wrap
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
        return txBox

    def _shape_text(
        self, shape, text: str,
        size: float = 12,
        color: RGBColor = WHITE,
        font: str = "Quicksand",
        bold: bool = False,
        align: PP_ALIGN = PP_ALIGN.LEFT,
    ):
        """Add text directly to a shape's text frame."""
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color

    def _image(self, slide, rel_path: str, x, y, w, h):
        """Add image from the Potomac assets folder (silently skips if missing)."""
        full = self.images / rel_path
        if not full.exists():
            logger.warning(f"Asset not found: {full}")
            return
        try:
            slide.shapes.add_picture(str(full), Inches(x), Inches(y), Inches(w), Inches(h))
        except Exception as e:
            logger.warning(f"Could not embed image {full}: {e}")

    def _logo(self, slide, filename: str, x, y, w, h):
        """Add a logo from the brand-assets/logos folder."""
        full = self.logos / filename
        if not full.exists():
            return
        try:
            slide.shapes.add_picture(str(full), Inches(x), Inches(y), Inches(w), Inches(h))
        except Exception as e:
            logger.warning(f"Could not embed logo {full}: {e}")

    # ─── brand layout helpers ────────────────────────────────────────────────

    def _add_bg(self, slide, color: RGBColor):
        self._rect(slide, 0, 0, SLIDE_W, SLIDE_H, color)

    def _add_dark_header(self, slide, title: str):
        """Dark header bar with yellow text and left yellow stripe."""
        self._rect(slide, 0, 0, SLIDE_W, 1.0, DARK)
        self._rect(slide, 0, 0, 0.12, 1.0, YELLOW)
        self._text(slide, title.upper(), 0.22, 0.08, 12.0, 0.82,
                   size=28, color=YELLOW, font="Rajdhani", bold=True, align=PP_ALIGN.LEFT)

    def _add_yellow_accent(self, slide):
        """Thin yellow left bar for content slides."""
        self._rect(slide, 0, 0, 0.12, SLIDE_H, YELLOW)

    def _add_teal_rule(self, slide, y: float = 1.15, w: float = 12.9):
        self._rect(slide, 0.22, y, w, 0.04, TEAL)

    def _add_footer(self, slide, dark: bool = False):
        text_color = WHITE if dark else MID_GRAY
        self._text(slide,
                   "Built to Conquer Risk®  |  potomacfundmanagement.com",
                   0.22, 7.15, 11.0, 0.25,
                   size=8, color=text_color, font="Quicksand", italic=True)
        self._text(slide, "CONFIDENTIAL",
                   11.5, 7.15, 1.6, 0.25,
                   size=8, color=YELLOW, font="Quicksand", bold=True, align=PP_ALIGN.RIGHT)

    def _dark_panel(self, slide, x: float = 8.9, y: float = 1.1, w: float = 4.2, h: float = 5.9):
        """Add a dark right-side info panel, returns (x, y) origin."""
        self._rect(slide, x, y, w, h, DARK)
        return x, y

    def _panel_header(self, slide, title: str, x: float, y: float, color: RGBColor = YELLOW):
        self._text(slide, title, x + 0.1, y + 0.15, 3.9, 0.4,
                   size=13, color=color, font="Rajdhani", bold=True)
        self._rect(slide, x + 0.1, y + 0.55, 3.9, 0.03, color)

    # ─── slide builders ──────────────────────────────────────────────────────

    def _slide_title(self, prs, data: dict):
        """Slide 1 — Dark title slide with yellow panel and tagline."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        self._add_bg(slide, DARK)

        # Yellow left panel
        self._rect(slide, 0, 0, 0.5, SLIDE_H, YELLOW)
        # Yellow horizontal accent
        self._rect(slide, 0.5, 3.5, 12.83, 0.06, YELLOW)
        # Teal accent above title
        self._rect(slide, 0.5, 2.1, 6.0, 0.05, TEAL)

        title    = data.get("title", "MARKET OUTLOOK").upper()
        subtitle = data.get("subtitle", "Q1 2025 — ECONOMIC & INVESTMENT REVIEW").upper()
        date     = data.get("date", "2025")
        presenter = data.get("presenter", "Investment Strategy & Research Team")

        self._text(slide, title, 0.7, 2.25, 10.0, 0.85,
                   size=48, color=WHITE, font="Rajdhani", bold=True)
        self._text(slide, subtitle, 0.7, 3.15, 10.0, 0.45,
                   size=20, color=YELLOW, font="Rajdhani")
        self._text(slide, date, 0.7, 3.68, 5.0, 0.35,
                   size=14, color=GRAY_999, font="Quicksand")
        self._text(slide, "Potomac Fund Management", 0.7, 4.4, 8.0, 0.35,
                   size=16, color=WHITE, font="Quicksand", bold=True)
        self._text(slide, presenter, 0.7, 4.75, 8.0, 0.3,
                   size=13, color=GRAY_999, font="Quicksand")

        self._logo(slide, "potomac-icon-yellow.png", 11.0, 6.5, 1.8, 0.75)
        self._add_footer(slide, dark=True)

    def _slide_agenda(self, prs, data: dict):
        """Slide 2 — Agenda with numbered topic grid."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_bg(slide, LIGHT_BG)
        self._add_yellow_accent(slide)
        self._add_dark_header(slide, "AGENDA")

        topics = data.get("topics", [])
        for i, t in enumerate(topics[:6]):
            col = 0 if i < 3 else 1
            row = i % 3
            x = 0.4 if col == 0 else 6.9
            y = 1.35 + row * 1.85

            # Numbered circle
            circ = self._rect(slide, x, y + 0.1, 0.55, 0.55, YELLOW, OVAL)
            self._shape_text(circ, t.get("num", f"0{i+1}"),
                             size=14, color=DARK, font="Rajdhani", bold=True, align=PP_ALIGN.CENTER)

            self._text(slide, t.get("title", "").upper(), x + 0.65, y, 5.5, 0.35,
                       size=15, color=DARK, font="Rajdhani", bold=True)
            self._text(slide, t.get("sub", ""), x + 0.65, y + 0.35, 5.5, 0.5,
                       size=10.5, color=MID_GRAY, font="Quicksand")
            line_color = TEAL if col == 0 else YELLOW
            self._rect(slide, x, y + 0.85, 6.0, 0.02, line_color)

        self._logo(slide, "potomac-icon-black.png", 11.8, 6.8, 1.3, 0.55)
        self._add_footer(slide)

    def _slide_chart(self, prs, data: dict):
        """Chart slide with left image + right dark callout panel."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        bg_color = LIGHT_BG if data.get("background") == "light" else WHITE
        self._add_bg(slide, bg_color)
        self._add_dark_header(slide, data.get("title", "ANALYSIS"))
        if bg_color == LIGHT_BG:
            self._add_yellow_accent(slide)

        # Chart image
        image_file = data.get("image", "")
        if image_file:
            self._image(slide, image_file, 0.2, 1.1, 8.5, 5.9)

        # Right dark panel
        px, py = self._dark_panel(slide)
        panel_title = data.get("panel_title", "KEY TAKEAWAYS")
        panel_color_name = data.get("panel_color", "yellow")
        panel_color = TEAL if panel_color_name == "teal" else YELLOW
        self._panel_header(slide, panel_title, px, py, panel_color)

        callouts = data.get("callouts", [])
        for i, c in enumerate(callouts[:5]):
            y = py + 0.7 + i * 1.0
            label = c.get("label", "").upper()
            text  = c.get("text", "")
            # Label badge
            badge = self._rect(slide, px + 0.1, y, 0.9, 0.28, YELLOW)
            self._shape_text(badge, label, size=8, color=DARK, font="Rajdhani",
                             bold=True, align=PP_ALIGN.CENTER)
            self._text(slide, text, px + 1.1, y, 2.9, 0.75,
                       size=10, color=WHITE, font="Quicksand")

        # Bullet points as fallback if no callouts
        bullets = data.get("bullets", [])
        if not callouts and bullets:
            for i, b in enumerate(bullets[:6]):
                y = py + 0.7 + i * 0.85
                self._rect(slide, px + 0.2, y + 0.12, 0.06, 0.06, TEAL)
                self._text(slide, b, px + 0.35, y, 3.6, 0.75,
                           size=10, color=WHITE, font="Quicksand")

        self._add_footer(slide)

    def _slide_stats(self, prs, data: dict):
        """Stats panel slide: chart image + stat values on dark panel."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_bg(slide, LIGHT_BG)
        self._add_dark_header(slide, data.get("title", "KEY STATISTICS"))
        self._add_yellow_accent(slide)

        image_file = data.get("image", "")
        if image_file:
            self._image(slide, image_file, 0.2, 1.1, 8.5, 5.9)

        px, py = self._dark_panel(slide)
        panel_title = data.get("panel_title", "SNAPSHOT")
        self._panel_header(slide, panel_title, px, py, YELLOW)

        stats = data.get("stats", [])
        stat_colors = [YELLOW, TEAL, WHITE, YELLOW]
        for i, s in enumerate(stats[:4]):
            y = py + 0.75 + i * 1.25
            val = s.get("val", "—")
            label = s.get("label", "").upper()
            color = stat_colors[i % len(stat_colors)]
            self._text(slide, val, px + 0.1, y, 3.9, 0.6,
                       size=34, color=color, font="Rajdhani", bold=True, align=PP_ALIGN.CENTER)
            self._text(slide, label, px + 0.1, y + 0.58, 3.9, 0.25,
                       size=10, color=GRAY_999, font="Quicksand", align=PP_ALIGN.CENTER)
            if i < len(stats) - 1:
                self._rect(slide, px + 0.5, y + 0.9, 2.9, 0.01, RGBColor(0x44, 0x44, 0x44))

        self._add_footer(slide)

    def _slide_two_charts(self, prs, data: dict):
        """Two side-by-side chart images with labeled footers."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_bg(slide, LIGHT_BG)
        self._add_dark_header(slide, data.get("title", "COMPARATIVE ANALYSIS"))
        self._add_yellow_accent(slide)

        charts = data.get("charts", [])
        positions = [(0.2, 6.3), (6.8, 6.3)]
        label_colors = [YELLOW, TEAL]

        for i, chart in enumerate(charts[:2]):
            x, w = positions[i]
            image_file = chart.get("image", "")
            if image_file:
                self._image(slide, image_file, x, 1.1, w, 5.5)
            # Label bar below chart
            lbl_color = label_colors[i % 2]
            lbl = self._rect(slide, x, 6.55, w, 0.3, lbl_color)
            self._shape_text(lbl, chart.get("label", "").upper(),
                             size=11, color=DARK, font="Rajdhani", bold=True, align=PP_ALIGN.CENTER)

        self._add_footer(slide)

    def _slide_content(self, prs, data: dict):
        """Generic content slide with bullet points (light or dark background)."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        dark = data.get("background") == "dark"
        bg_color = DARK if dark else LIGHT_BG
        self._add_bg(slide, bg_color)
        self._add_dark_header(slide, data.get("title", ""))

        bullets = data.get("bullets", [])
        text_color = WHITE if dark else DARK
        for i, b in enumerate(bullets[:8]):
            y = 1.25 + i * 0.72
            dot_color = YELLOW if dark else TEAL
            self._rect(slide, 0.4, y + 0.12, 0.08, 0.08, dot_color)
            self._text(slide, b, 0.6, y, 12.3, 0.65,
                       size=13, color=text_color, font="Quicksand")

        self._add_footer(slide, dark=dark)

    def _slide_summary(self, prs, data: dict):
        """3-column portfolio summary (Overweight / Neutral / Underweight)."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_bg(slide, WHITE)
        self._add_dark_header(slide, data.get("title", "POSITIONING SUMMARY"))

        columns = data.get("columns", [])
        col_defaults = [
            {"title": "OVERWEIGHT", "color": "green", "items": []},
            {"title": "NEUTRAL",    "color": "yellow", "items": []},
            {"title": "UNDERWEIGHT","color": "red",    "items": []},
        ]
        if not columns:
            columns = col_defaults

        color_map = {"green": GREEN, "yellow": YELLOW, "red": RED, "teal": TEAL}

        for ci, col in enumerate(columns[:3]):
            x = 0.3 + ci * 4.33
            col_color = color_map.get(col.get("color", "yellow"), YELLOW)

            # Header bar
            hdr = self._rect(slide, x, 1.1, 4.0, 0.5, col_color)
            self._shape_text(hdr, col.get("title", "").upper(),
                             size=15, color=DARK, font="Rajdhani", bold=True, align=PP_ALIGN.CENTER)

            # Items
            items = col.get("items", [])
            for ii, item in enumerate(items[:5]):
                iy = 1.75 + ii * 0.95
                row_bg = LIGHT_BG if ii % 2 == 0 else WHITE
                row = slide.shapes.add_shape(RECT, Inches(x), Inches(iy), Inches(4.0), Inches(0.75))
                row.fill.solid()
                row.fill.fore_color.rgb = row_bg
                row.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
                # Color left accent on row
                self._rect(slide, x, iy, 0.08, 0.75, col_color)
                self._text(slide, item, x + 0.15, iy, 3.75, 0.75,
                           size=12, color=DARK, font="Quicksand")

        self._add_footer(slide)

    def _slide_closing(self, prs, data: dict):
        """Closing slide — full dark background with tagline and contact grid."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_bg(slide, DARK)
        # Yellow top & bottom stripes
        self._rect(slide, 0, 0, SLIDE_W, 0.12, YELLOW)
        self._rect(slide, 0, 7.38, SLIDE_W, 0.12, YELLOW)

        tagline = data.get("tagline", "BUILT TO CONQUER RISK®").upper()
        self._text(slide, tagline, 0.5, 1.8, 12.33, 1.0,
                   size=52, color=YELLOW, font="Rajdhani", bold=True, align=PP_ALIGN.CENTER)
        self._rect(slide, 4.0, 2.9, 5.33, 0.05, TEAL)
        self._text(slide, "Systematic. Disciplined. Risk-Managed.",
                   0.5, 3.05, 12.33, 0.5,
                   size=20, color=WHITE, font="Quicksand", italic=True, align=PP_ALIGN.CENTER)

        contacts = data.get("contacts", [
            {"label": "WEBSITE",  "val": "potomacfundmanagement.com"},
            {"label": "EMAIL",    "val": "info@potomacfm.com"},
            {"label": "PHONE",    "val": "(800) 851-0933"},
            {"label": "LOCATION", "val": "Bethesda, MD"},
        ])
        for i, c in enumerate(contacts[:4]):
            x = 0.5 + i * 3.1
            self._text(slide, c.get("label", "").upper(), x, 4.5, 2.8, 0.3,
                       size=9, color=YELLOW, font="Rajdhani", bold=True, align=PP_ALIGN.CENTER)
            self._text(slide, c.get("val", ""), x, 4.82, 2.8, 0.35,
                       size=11, color=WHITE, font="Quicksand", align=PP_ALIGN.CENTER)

        disclaimer = data.get("disclaimer",
            "This presentation is for informational purposes only and does not constitute investment advice. "
            "Past performance is not indicative of future results. Investing involves risk, including loss of principal. "
            "Potomac Fund Management is a registered investment adviser.")
        self._text(slide, disclaimer, 0.5, 6.3, 12.33, 0.65,
                   size=7.5, color=GRAY_999, font="Quicksand", italic=True, align=PP_ALIGN.CENTER)

        self._logo(slide, "potomac-icon-yellow.png", 5.5, 5.5, 2.3, 0.65)

    # ─── slide dispatcher ────────────────────────────────────────────────────

    def _build_slide(self, prs, slide_data: dict):
        """Route to the correct slide builder based on 'type'."""
        slide_type = slide_data.get("type", "content")
        dispatch = {
            "title":       self._slide_title,
            "agenda":      self._slide_agenda,
            "chart":       self._slide_chart,
            "stats":       self._slide_stats,
            "two_charts":  self._slide_two_charts,
            "content":     self._slide_content,
            "summary":     self._slide_summary,
            "closing":     self._slide_closing,
        }
        builder = dispatch.get(slide_type, self._slide_content)
        builder(prs, slide_data)

    # ─── public API ──────────────────────────────────────────────────────────

    def generate_from_outline(self, outline: dict, output_path: str) -> str:
        """
        Build a Potomac .pptx from an AI-generated outline dict.

        outline format:
        {
          "title": "Presentation Title",
          "slides": [
            {"type": "title",  "title": "...", "subtitle": "...", "date": "..."},
            {"type": "agenda", "topics": [{"num": "01", "title": "...", "sub": "..."}]},
            {"type": "chart",  "title": "...", "image": "filename.png", "callouts": [...]},
            {"type": "stats",  "title": "...", "image": "...", "stats": [{"val":"3.4%","label":"CPI"}]},
            {"type": "two_charts", "title": "...", "charts": [{"image":"...","label":"..."}]},
            {"type": "content","title": "...", "bullets": ["..."]},
            {"type": "summary","title": "...", "columns": [{"title":"OW","color":"green","items":[...]}]},
            {"type": "closing","tagline": "...", "contacts": [...]},
          ]
        }
        """
        prs = Presentation()
        prs.slide_width  = Inches(SLIDE_W)
        prs.slide_height = Inches(SLIDE_H)

        # Ensure at least a blank layout exists
        if not prs.slide_layouts:
            raise RuntimeError("python-pptx: no slide layouts found in blank prs")

        slides = outline.get("slides", [])
        if not slides:
            # Minimal fallback
            slides = [
                {"type": "title",   "title": outline.get("title", "Presentation"),
                 "subtitle": outline.get("subtitle", ""), "date": "2025"},
                {"type": "closing"},
            ]

        for slide_data in slides:
            try:
                self._build_slide(prs, slide_data)
            except Exception as e:
                logger.error(f"Error building slide {slide_data.get('type')}: {e}", exc_info=True)

        out = Path(output_path)
        out.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(out))
        logger.info(f"Saved presentation: {out} ({len(slides)} slides)")
        return str(out)

    def generate_market_outlook(
        self,
        output_path: str,
        title: str = "MARKET OUTLOOK",
        subtitle: str = "Q1 2025 — ECONOMIC & INVESTMENT REVIEW",
        date: str = "February 2025",
    ) -> str:
        """Convenience wrapper that generates the standard Market Outlook deck."""
        outline = {
            "title": title,
            "slides": [
                {
                    "type": "title",
                    "title": title,
                    "subtitle": subtitle,
                    "date": date,
                },
                {
                    "type": "agenda",
                    "topics": [
                        {"num": "01", "title": "EQUITY MARKETS",
                         "sub": "S&P 500 performance, momentum signals & sector rotation"},
                        {"num": "02", "title": "INFLATION & CPI",
                         "sub": "Consumer price trends, Fed policy implications"},
                        {"num": "03", "title": "FIXED INCOME",
                         "sub": "Credit spreads, duration risk, LQD performance"},
                        {"num": "04", "title": "COMMODITIES & GOLD",
                         "sub": "GLD, MOVE index, macro hedge positioning"},
                        {"num": "05", "title": "MARKET SCORING",
                         "sub": "Potomac risk scoring basket & tactical positioning"},
                        {"num": "06", "title": "PASSIVE vs. ACTIVE",
                         "sub": "Performance comparison & AFG fee impact analysis"},
                    ],
                },
                {
                    "type": "chart",
                    "title": "EQUITY MARKETS — S&P 500",
                    "image": "1 - SPX.png",
                    "panel_title": "KEY TAKEAWAYS",
                    "callouts": [
                        {"label": "TREND",  "text": "SPX momentum remains constructive above 200-DMA"},
                        {"label": "SIGNAL", "text": "Bull signal triggered Q4 2024; holding through correction"},
                        {"label": "RISK",   "text": "Watch 4,800 support — break below triggers defensive shift"},
                        {"label": "SECTOR", "text": "Tech & Energy leading; Utilities & Staples lagging"},
                    ],
                },
                {
                    "type": "stats",
                    "title": "INFLATION & CPI TRENDS",
                    "image": "Slide 7 - CPI.png",
                    "panel_title": "INFLATION SNAPSHOT",
                    "stats": [
                        {"val": "3.4%",    "label": "HEADLINE CPI"},
                        {"val": "3.9%",    "label": "CORE CPI"},
                        {"val": "5.25%",   "label": "FED FUNDS RATE"},
                        {"val": "Q2 2025", "label": "FIRST CUT EST."},
                    ],
                    "background": "light",
                },
                {
                    "type": "chart",
                    "title": "FIXED INCOME — CREDIT MARKETS",
                    "image": "4 - LQD.png",
                    "panel_title": "FIXED INCOME VIEW",
                    "panel_color": "teal",
                    "bullets": [
                        "IG credit spreads at 120 bps — approaching tight historical range",
                        "Duration risk elevated as Fed rate cuts get repriced lower",
                        "LQD showing negative momentum; short-duration tilt preferred",
                        "HY spreads at 350 bps — not pricing default risk adequately",
                        "Floating rate & T-Bills remain attractive vs. long IG bonds",
                    ],
                },
                {
                    "type": "two_charts",
                    "title": "COMMODITIES — GOLD & VOLATILITY",
                    "charts": [
                        {"image": "GLD.png",  "label": "GLD — GOLD ETF PERFORMANCE"},
                        {"image": "MOVE.png", "label": "MOVE INDEX — BOND MARKET VOLATILITY"},
                    ],
                    "background": "light",
                },
                {
                    "type": "chart",
                    "title": "POTOMAC RISK SCORING BASKET",
                    "image": "Scoring_basket.png",
                    "panel_title": "CURRENT SIGNALS",
                    "callouts": [
                        {"label": "SPX",     "text": "BULL — momentum above 200-DMA"},
                        {"label": "CREDIT",  "text": "NEUTRAL — spreads compressed"},
                        {"label": "VOL",     "text": "BEAR — MOVE index elevated"},
                        {"label": "MOMENTUM","text": "BULL — positive breadth trend"},
                        {"label": "BREADTH", "text": "NEUTRAL — mixed sector signals"},
                    ],
                },
                {
                    "type": "two_charts",
                    "title": "PASSIVE vs. ACTIVE — PERFORMANCE ANALYSIS",
                    "charts": [
                        {"image": "Market_Outlook.002.jpeg", "label": "PASSIVE vs. ACTIVE TABLE"},
                        {"image": "Market_Outlook.026.jpeg", "label": "AFG FEE IMPACT ANALYSIS"},
                    ],
                    "background": "light",
                },
                {
                    "type": "summary",
                    "title": "PORTFOLIO POSITIONING SUMMARY",
                    "columns": [
                        {
                            "title": "OVERWEIGHT",
                            "color": "green",
                            "items": ["US Large Cap Equities", "Short Duration Fixed Income",
                                      "Gold / Commodities", "Cash & T-Bills", "Energy Sector"],
                        },
                        {
                            "title": "NEUTRAL",
                            "color": "yellow",
                            "items": ["International Developed", "Investment Grade Credit",
                                      "Healthcare Sector", "Mid Cap Equities", "Real Estate (REITs)"],
                        },
                        {
                            "title": "UNDERWEIGHT",
                            "color": "red",
                            "items": ["Long Duration Bonds", "Emerging Markets",
                                      "High Yield Credit", "Utilities", "Consumer Discretionary"],
                        },
                    ],
                },
                {"type": "closing"},
            ],
        }
        return self.generate_from_outline(outline, output_path)
