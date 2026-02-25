"""
Supabase Storage Helper
=======================
Provides a unified interface for file operations with Supabase Storage.

Key Features:
- Magic byte validation for file types (security)
- File size limits per bucket
- Content hash deduplication
- File metadata tracking in database
- Automatic text extraction for indexing
"""

import hashlib
import mimetypes
import os
import uuid
import logging
from typing import Optional, Dict, Any, Tuple, List
from supabase import Client

logger = logging.getLogger(__name__)

# Allowed MIME types per bucket
ALLOWED_MIME_TYPES = {
    "user-uploads": [
        "application/pdf",
        "text/plain",
        "text/csv",
        "application/json",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        "image/png",
        "image/jpeg",
        "image/gif",
        "image/webp",
    ],
    "presentations": [
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    ],
    "brain-docs": None,  # Allow all types
}

# Max file sizes per bucket (in bytes)
MAX_FILE_SIZES = {
    "user-uploads": 50 * 1024 * 1024,  # 50MB
    "presentations": 100 * 1024 * 1024,  # 100MB
    "brain-docs": 50 * 1024 * 1024,  # 50MB
}

# Magic bytes for file type validation (first 4-8 bytes)
MAGIC_BYTES = {
    b'\x25\x50\x44\x46': 'application/pdf',  # PDF
    b'\x50\x4B\x03\x04': 'application/zip',  # ZIP (includes .pptx, .xlsx, .docx)
    b'\x89\x50\x4E\x47': 'image/png',  # PNG
    b'\xFF\xD8\xFF': 'image/jpeg',  # JPEG
    b'\x47\x49\x46\x38': 'image/gif',  # GIF
    b'\x52\x49\x46\x46': 'image/webp',  # WEBP (also WAV, AVI - check further)
}


class StorageError(Exception):
    """Custom exception for storage operations."""
    pass


class FileTypeValidationError(StorageError):
    """Raised when file type validation fails."""
    pass


class FileSizeExceededError(StorageError):
    """Raised when file size exceeds limit."""
    pass


class StorageHelper:
    """
    Manages file operations with Supabase Storage.
    
    Usage:
        storage = StorageHelper(supabase_client)
        
        # Upload a file
        file_record = await storage.upload_file(
            user_id="user-uuid",
            bucket="user-uploads",
            filename="document.pdf",
            content=file_bytes
        )
        
        # Get file for download
        file_info = await storage.get_file(file_id, user_id)
        
        # Delete a file
        deleted = await storage.delete_file(file_id, user_id)
    """
    
    def __init__(self, client: Client):
        """
        Initialize storage helper.
        
        Args:
            client: Supabase client instance (service_role recommended for backend)
        """
        self.client = client
        self._pending_uploads: Dict[str, Dict] = {}  # For presigned URL flow
    
    def _get_storage_path(self, user_id: str, filename: str) -> str:
        """
        Generate storage path for a file.
        
        Format: {user_id}/{uuid}_{sanitized_filename}
        """
        # Sanitize filename - remove dangerous characters
        safe_filename = "".join(
            c for c in filename if c.isalnum() or c in "._- "
        ).rstrip(".")[:100]
        
        # Replace spaces with underscores
        safe_filename = safe_filename.replace(" ", "_")
        
        # Add short UUID to prevent collisions
        file_uuid = str(uuid.uuid4())[:8]
        
        return f"{user_id}/{file_uuid}_{safe_filename}"
    
    def _compute_hash(self, content: bytes) -> str:
        """Compute SHA-256 hash of content for deduplication."""
        return hashlib.sha256(content).hexdigest()
    
    def _detect_mime_type(self, content: bytes, filename: str) -> str:
        """
        Detect MIME type from magic bytes (more secure than extension).
        
        Falls back to extension-based detection if magic bytes don't match.
        """
        # Try magic bytes first (more secure)
        for magic, mime in MAGIC_BYTES.items():
            if content[:len(magic)] == magic:
                # Special handling for Office files (ZIP-based containers)
                if mime == 'application/zip':
                    ext = os.path.splitext(filename)[1].lower()
                    office_types = {
                        '.pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
                        '.xlsx': 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
                        '.docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
                    }
                    if ext in office_types:
                        return office_types[ext]
                return mime
        
        # Fallback to mimetypes library (extension-based)
        mime_type, _ = mimetypes.guess_type(filename)
        return mime_type or 'application/octet-stream'
    
    def validate_file(
        self, 
        content: bytes, 
        filename: str, 
        bucket: str
    ) -> Tuple[bool, Optional[str]]:
        """
        Validate file type and size.
        
        Args:
            content: File content bytes
            filename: Original filename
            bucket: Target bucket ID
            
        Returns:
            Tuple of (is_valid, error_message)
        """
        # Check file size
        max_size = MAX_FILE_SIZES.get(bucket, 50 * 1024 * 1024)
        if len(content) > max_size:
            size_mb = len(content) / (1024 * 1024)
            max_mb = max_size / (1024 * 1024)
            return False, f"File size ({size_mb:.1f}MB) exceeds limit ({max_mb:.0f}MB) for bucket {bucket}"
        
        # Check MIME type
        detected_type = self._detect_mime_type(content, filename)
        allowed = ALLOWED_MIME_TYPES.get(bucket)
        
        if allowed and detected_type not in allowed:
            return False, f"File type '{detected_type}' not allowed in bucket '{bucket}'"
        
        return True, None
    
    async def get_upload_url(
        self,
        user_id: str,
        bucket: str,
        filename: str,
        content_type: Optional[str] = None,
        expires_in: int = 3600,
    ) -> Dict[str, Any]:
        """
        Generate upload info for client-side direct upload.
        
        Since Supabase doesn't have native presigned URLs like S3,
        this returns the info needed for the client to upload via
        the Supabase JS SDK.
        
        Args:
            user_id: User UUID
            bucket: Bucket ID (user-uploads, presentations, brain-docs)
            filename: Original filename
            content_type: MIME type (optional)
            expires_in: Not used (for S3 compatibility)
            
        Returns:
            Dict with upload instructions
        """
        if bucket not in MAX_FILE_SIZES:
            raise StorageError(f"Invalid bucket: {bucket}")
        
        storage_path = self._get_storage_path(user_id, filename)
        upload_id = str(uuid.uuid4())
        
        # Store pending upload info (for confirmation later)
        self._pending_uploads[upload_id] = {
            "user_id": user_id,
            "bucket": bucket,
            "storage_path": storage_path,
            "filename": filename,
            "content_type": content_type,
            "created_at": os.times().elapsed,
        }
        
        return {
            "upload_id": upload_id,
            "bucket": bucket,
            "storage_path": storage_path,
            "filename": filename,
            "max_size": MAX_FILE_SIZES.get(bucket),
            "instructions": {
                "method": "POST",
                "content_type": content_type or "application/octet-stream",
            },
        }
    
    async def upload_file(
        self,
        user_id: str,
        bucket: str,
        filename: str,
        content: bytes,
        content_type: Optional[str] = None,
    ) -> Dict[str, Any]:
        """
        Upload a file directly (server-side upload).
        
        Args:
            user_id: User UUID
            bucket: Bucket ID
            filename: Original filename
            content: File content bytes
            content_type: MIME type (optional, auto-detected)
            
        Returns:
            File record dict with id, storage_path, etc.
        """
        # Validate file
        is_valid, error = self.validate_file(content, filename, bucket)
        if not is_valid:
            raise FileTypeValidationError(error)
        
        # Detect content type from magic bytes
        if not content_type:
            content_type = self._detect_mime_type(content, filename)
        
        # Generate storage path
        storage_path = self._get_storage_path(user_id, filename)
        
        # Compute hash for deduplication
        content_hash = self._compute_hash(content)
        
        # Check for existing file with same hash (deduplication)
        existing = self.client.table("file_uploads").select("*").eq(
            "user_id", user_id
        ).eq("content_hash", content_hash).execute()
        
        if existing.data:
            logger.info(f"Duplicate file detected: {content_hash[:16]}...")
            return existing.data[0]
        
        # Upload to Supabase Storage
        try:
            self.client.storage.from_(bucket).upload(
                path=storage_path,
                file=content,
                file_options={
                    "content-type": content_type,
                    "upsert": "true",
                }
            )
        except Exception as e:
            error_msg = str(e)
            logger.error(f"Storage upload failed: {error_msg}")
            raise StorageError(f"Failed to upload file: {error_msg}")
        
        # Create database record
        file_record = {
            "user_id": user_id,
            "bucket_id": bucket,
            "storage_path": storage_path,
            "original_filename": filename,
            "content_type": content_type,
            "file_size": len(content),
            "content_hash": content_hash,
            "status": "uploaded",
        }
        
        result = self.client.table("file_uploads").insert(file_record).execute()
        
        if not result.data:
            raise StorageError("Failed to create file record")
        
        logger.info(f"File uploaded: {storage_path} ({len(content)} bytes)")
        return result.data[0]
    
    async def get_file(
        self,
        file_id: str,
        user_id: str,
    ) -> Optional[Dict[str, Any]]:
        """
        Get file metadata and download URL.
        
        Args:
            file_id: File upload UUID
            user_id: User UUID (for ownership check)
            
        Returns:
            File record with download_url, or None if not found
        """
        result = self.client.table("file_uploads").select("*").eq(
            "id", file_id
        ).eq("user_id", user_id).execute()
        
        if not result.data:
            return None
        
        file_record = result.data[0]
        
        # Generate signed URL for download (valid for 1 hour)
        try:
            signed_url = self.client.storage.from_(
                file_record["bucket_id"]
            ).create_signed_url(
                file_record["storage_path"],
                expires_in=3600
            )
            file_record["download_url"] = signed_url.get("signedURL")
        except Exception as e:
            logger.warning(f"Failed to generate signed URL: {e}")
            file_record["download_url"] = None
        
        return file_record
    
    async def get_file_content(
        self,
        file_id: str,
        user_id: str,
    ) -> Optional[bytes]:
        """
        Download file content directly.
        
        Args:
            file_id: File upload UUID
            user_id: User UUID
            
        Returns:
            File content bytes, or None if not found
        """
        file_record = await self.get_file(file_id, user_id)
        if not file_record:
            return None
        
        try:
            response = self.client.storage.from_(
                file_record["bucket_id"]
            ).download(file_record["storage_path"])
            return response
        except Exception as e:
            logger.error(f"Failed to download file: {e}")
            return None
    
    async def delete_file(
        self,
        file_id: str,
        user_id: str,
    ) -> bool:
        """
        Delete a file from storage and database.
        
        Args:
            file_id: File upload UUID
            user_id: User UUID
            
        Returns:
            True if deleted, False if not found
        """
        # Get file record
        result = self.client.table("file_uploads").select("*").eq(
            "id", file_id
        ).eq("user_id", user_id).execute()
        
        if not result.data:
            return False
        
        file_record = result.data[0]
        
        # Delete from storage
        try:
            self.client.storage.from_(
                file_record["bucket_id"]
            ).remove([file_record["storage_path"]])
        except Exception as e:
            logger.warning(f"Failed to delete from storage: {e}")
            # Continue to delete database record even if storage delete fails
        
        # Delete from database
        self.client.table("file_uploads").delete().eq("id", file_id).execute()
        
        logger.info(f"File deleted: {file_id}")
        return True
    
    async def list_files(
        self,
        user_id: str,
        bucket: Optional[str] = None,
        limit: int = 50,
        offset: int = 0,
    ) -> List[Dict[str, Any]]:
        """
        List files for a user.
        
        Args:
            user_id: User UUID
            bucket: Optional bucket filter
            limit: Max results
            offset: Pagination offset
            
        Returns:
            List of file records
        """
        query = self.client.table("file_uploads").select("*").eq(
            "user_id", user_id
        ).order("created_at", desc=True).limit(limit).offset(offset)
        
        if bucket:
            query = query.eq("bucket_id", bucket)
        
        result = query.execute()
        return result.data or []
    
    async def extract_text(
        self,
        file_id: str,
        user_id: str,
    ) -> Optional[str]:
        """
        Extract text content from a file.
        
        Supports: PDF, TXT, CSV, JSON, PPTX
        
        Args:
            file_id: File upload UUID
            user_id: User UUID
            
        Returns:
            Extracted text, or None if extraction failed
        """
        content = await self.get_file_content(file_id, user_id)
        if not content:
            return None
        
        file_record = await self.get_file(file_id, user_id)
        if not file_record:
            return None
        
        content_type = file_record.get("content_type", "")
        
        try:
            # Plain text files
            if content_type in ("text/plain", "text/csv", "application/json"):
                return content.decode("utf-8", errors="ignore")
            
            # PDF files
            elif content_type == "application/pdf":
                try:
                    import io
                    from pdfplumber import PDF
                    
                    pdf = PDF(io.BytesIO(content))
                    text = ""
                    for page in pdf.pages:
                        page_text = page.extract_text()
                        if page_text:
                            text += page_text + "\n"
                    return text.strip()
                except ImportError:
                    logger.warning("pdfplumber not installed, cannot extract PDF text")
                    return None
            
            # PowerPoint files
            elif "presentationml" in content_type:
                try:
                    from pptx import Presentation
                    import io
                    
                    prs = Presentation(io.BytesIO(content))
                    text = ""
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                text += shape.text + "\n"
                    return text.strip()
                except ImportError:
                    logger.warning("python-pptx not installed, cannot extract PPTX text")
                    return None
            
            # Excel files
            elif "spreadsheetml" in content_type:
                try:
                    import pandas as pd
                    import io
                    
                    df = pd.read_excel(io.BytesIO(content))
                    return df.to_string()
                except ImportError:
                    logger.warning("pandas not installed, cannot extract Excel text")
                    return None
            
            else:
                logger.warning(f"Unsupported content type for text extraction: {content_type}")
                return None
                
        except Exception as e:
            logger.error(f"Text extraction failed: {e}")
            return None
    
    async def confirm_upload(
        self,
        upload_id: str,
        user_id: str,
        file_size: Optional[int] = None,
    ) -> Dict[str, Any]:
        """
        Confirm a client-side upload and create database record.
        
        Call this after client uploads directly via Supabase SDK.
        
        Args:
            upload_id: Upload ID from get_upload_url()
            user_id: User UUID
            file_size: Actual file size (optional)
            
        Returns:
            File record dict
        """
        pending = self._pending_uploads.get(upload_id)
        if not pending:
            raise StorageError("Upload session not found or expired")
        
        if pending["user_id"] != user_id:
            raise StorageError("Upload session belongs to another user")
        
        # Verify file exists in storage
        try:
            # Try to get file info from storage
            files = self.client.storage.from_(pending["bucket"]).list(
                path=pending["user_id"]
            )
            
            # Check if our file is there
            storage_filename = os.path.basename(pending["storage_path"])
            file_exists = any(
                f.get("name") == storage_filename for f in files
            )
            
            if not file_exists:
                raise StorageError("File not found in storage - upload may have failed")
                
        except StorageError:
            raise
        except Exception as e:
            logger.warning(f"Could not verify file in storage: {e}")
            # Continue anyway - the storage API might not support listing
        
        # Create database record
        db_record = {
            "user_id": user_id,
            "bucket_id": pending["bucket"],
            "storage_path": pending["storage_path"],
            "original_filename": pending["filename"],
            "content_type": pending.get("content_type"),
            "file_size": file_size,
            "status": "uploaded",
        }
        
        result = self.client.table("file_uploads").insert(db_record).execute()
        
        # Clean up pending upload
        del self._pending_uploads[upload_id]
        
        return result.data[0] if result.data else db_record


# Singleton helper instance
_storage_helper: Optional[StorageHelper] = None


def get_storage_helper(client: Optional[Client] = None) -> StorageHelper:
    """
    Get or create storage helper instance.
    
    Args:
        client: Optional Supabase client (uses default if not provided)
        
    Returns:
        StorageHelper instance
    """
    global _storage_helper
    
    if _storage_helper is None or client is not None:
        if client is None:
            from db.supabase_client import get_supabase
            client = get_supabase()
        _storage_helper = StorageHelper(client)
    
    return _storage_helper