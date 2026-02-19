"""
Potomac InDesign → PPTX Template Converter
==========================================
Reads the manifest JSON produced by export_manifest.jsx and builds
python-pptx template files for each slide type.

Usage:
    pip install python-pptx Pillow
    python convert_to_templates.py \
        --manifest ./potomac_export/Bull_Bear_manifest.json \
        --images   ./potomac_export/slide_images/ \
        --output   ./templates/bull-bear/ \
        --logos    ./brand_assets/logos/

The script does TWO things per slide:
  1. Bakes the full-slide reference PNG as a background image (preserves all
     complex vector/design elements perfectly)
  2. Overlays transparent python-pptx placeholders exactly where each
     editable text/image zone sits — so Claude can fill them at generation time.
"""

import json
import os
import sys
import argparse
import shutil
from pathlib import Path
from typing import Optional

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt
    import pptx.oxml.ns as pns
    from lxml import etree
except ImportError:
    print("Install required packages: pip install python-pptx lxml Pillow")
    sys.exit(1)

# ─── POTOMAC BRAND CONSTANTS ─────────────────────────────────────────────────

BRAND = {
    "yellow":    RGBColor(0xFE, 0xC0, 0x0F),
    "dark_gray": RGBColor(0x21, 0x21, 0x21),
    "turquoise": RGBColor(0x00, 0xDE, 0xD1),
    "pink":      RGBColor(0xEB, 0x2F, 0x5C),
    "white":     RGBColor(0xFF, 0xFF, 0xFF),
    "light_bg":  RGBColor(0xF2, 0xF2, 0xF2),
}

FONT_HEADER = "Rajdhani"
FONT_BODY   = "Quicksand"

# Standard Potomac widescreen: 13.333" x 7.5"
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ─── HELPERS ─────────────────────────────────────────────────────────────────

def hex_to_rgb(hex_str: str) -> Optional[RGBColor]:
    """Convert hex string (no #) to RGBColor."""
    if not hex_str:
        return None
    h = hex_str.lstrip("#")
    if len(h) != 6:
        return None
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def in_to_emu(inches: float) -> int:
    return int(inches * 914400)


def is_dark_background(hex_str: str) -> bool:
    """Determine if a hex color is dark (for choosing text color)."""
    if not hex_str:
        return False
    h = hex_str.lstrip("#")
    r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
    luminance = 0.299 * r + 0.587 * g + 0.114 * b
    return luminance < 128


def add_background_image(slide, image_path: str):
    """Set a PNG as the full-slide background."""
    if not image_path or not os.path.exists(image_path):
        return
    left = top = 0
    pic = slide.shapes.add_picture(image_path, left, top, SLIDE_W, SLIDE_H)
    # Move background image to back
    slide.shapes._spTree.remove(pic._element)
    slide.shapes._spTree.insert(2, pic._element)


def add_placeholder_textbox(slide, pos: dict, name: str, placeholder_text: str,
                             font: str = FONT_BODY, size_pt: float = 14,
                             bold: bool = False, all_caps: bool = False,
                             color: RGBColor = None, align: str = "left",
                             transparent: bool = True):
    """
    Add a named text placeholder over a specific zone.
    The box is transparent by default so the background image shows through.
    The placeholder text is there for python-pptx assembly to find & replace.
    """
    left   = Inches(pos["x_in"])
    top    = Inches(pos["y_in"])
    width  = Inches(pos["w_in"])
    height = Inches(pos["h_in"])

    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.name = name  # key: this is how the assembler finds it

    tf = txBox.text_frame
    tf.word_wrap = True

    # Set placeholder text
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = placeholder_text

    # Font
    run.font.name = font
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color

    # Alignment
    align_map = {
        "left":   PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right":  PP_ALIGN.RIGHT,
    }
    p.alignment = align_map.get(align.lower(), PP_ALIGN.LEFT)

    if all_caps:
        run.font.all_caps = True

    # Make fill transparent
    if transparent:
        fill = txBox.fill
        fill.background()

    return txBox


def add_placeholder_image_zone(slide, pos: dict, name: str):
    """
    Add a named empty rectangle as an image placeholder zone.
    The rectangle is transparent with a faint yellow border so it's
    visible during template debugging but won't show in final output.
    """
    left   = Inches(pos["x_in"])
    top    = Inches(pos["y_in"])
    width  = Inches(pos["w_in"])
    height = Inches(pos["h_in"])

    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.name = name

    # Transparent fill
    shape.fill.background()

    # Faint yellow border (debug only — remove for production)
    shape.line.color.rgb = BRAND["yellow"]
    shape.line.width = Pt(1)

    return shape


# ─── TEMPLATE BUILDERS ───────────────────────────────────────────────────────
# Each function takes a Presentation and a slide manifest dict, adds a slide,
# and registers all editable zones. The function name = template_id.

def build_template_cover_hero(prs: Presentation, slide_data: dict,
                               bg_image: str):
    """
    Slide 1: Cover hero
    - Full yellow bottom panel (~60% of height), white strip top
    - Logo top-left on white strip
    - Large left-aligned title on yellow area
    - Right side: decorative product image placeholder
    """
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)

    add_background_image(slide, bg_image)

    # Editable zones (positions calibrated from Bull Bear PDF)
    add_placeholder_textbox(slide,
        pos={"x_in": 0.4, "y_in": 2.5, "w_in": 6.5, "h_in": 2.8},
        name="main_title",
        placeholder_text="{{main_title}}",
        font=FONT_HEADER, size_pt=72, bold=True,
        all_caps=True, color=BRAND["dark_gray"], align="left"
    )
    add_placeholder_image_zone(slide,
        pos={"x_in": 6.8, "y_in": 1.2, "w_in": 6.1, "h_in": 5.8},
        name="hero_image"
    )
    return slide


def build_template_strategy_intro(prs: Presentation, slide_data: dict,
                                   bg_image: str):
    """
    Slide 2: Strategy intro / title card
    - Light gray background
    - Breadcrumb top-left, icon top-right
    - Large centered strategy name with icon image to the right of text
    - Centered subtitle paragraph below
    """
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    add_background_image(slide, bg_image)

    add_placeholder_textbox(slide,
        pos={"x_in": 0.35, "y_in": 0.18, "w_in": 6.0, "h_in": 0.32},
        name="breadcrumb",
        placeholder_text="{{breadcrumb}}",
        font=FONT_HEADER, size_pt=11, bold=False,
        all_caps=True, color=BRAND["dark_gray"], align="left"
    )
    add_placeholder_textbox(slide,
        pos={"x_in": 1.5, "y_in": 2.2, "w_in": 6.5, "h_in": 1.0},
        name="strategy_name",
        placeholder_text="{{strategy_name}}",
        font=FONT_HEADER, size_pt=64, bold=True,
        all_caps=True, color=BRAND["dark_gray"], align="left"
    )
    add_placeholder_image_zone(slide,
        pos={"x_in": 7.8, "y_in": 1.9, "w_in": 2.8, "h_in": 1.6},
        name="strategy_icon"
    )
    add_placeholder_textbox(slide,
        pos={"x_in": 1.5, "y_in": 3.9, "w_in": 10.3, "h_in": 1.2},
        name="subtitle",
        placeholder_text="{{subtitle}}",
        font=FONT_BODY, size_pt=20, bold=False,
        color=BRAND["dark_gray"], align="center"
    )
    return slide


def build_template_chart_full_light(prs: Presentation, slide_data: dict,
                                     bg_image: str):
    """
    Slides 5,6,7: Full-width Optuma chart slide
    - Light gray background
    - Small breadcrumb top-left ("PROCESS"), icon top-right
    - Large centered title (Rajdhani ALL CAPS)
    - Full-width chart image placeholder (fills most of slide)
    - "For illustrative purposes only." footer centered
    - Optional vertical sidebar text on left edge (very small)
    """
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    add_background_image(slide, bg_image)

    add_placeholder_textbox(slide,
        pos={"x_in": 0.35, "y_in": 0.18, "w_in": 3.0, "h_in": 0.32},
        name="breadcrumb",
        placeholder_text="{{breadcrumb}}",
        font=FONT_HEADER, size_pt=11, bold=False,
        all_caps=True, color=BRAND["dark_gray"], align="left"
    )
    add_placeholder_textbox(slide,
        pos={"x_in": 1.5, "y_in": 0.5, "w_in": 10.3, "h_in": 0.75},
        name="slide_title",
        placeholder_text="{{slide_title}}",
        font=FONT_HEADER, size_pt=40, bold=True,
        all_caps=True, color=BRAND["dark_gray"], align="center"
    )
    add_placeholder_image_zone(slide,
        pos={"x_in": 0.2, "y_in": 1.35, "w_in": 12.9, "h_in": 5.55},
        name="chart_image"
    )
    add_placeholder_textbox(slide,
        pos={"x_in": 1.5, "y_in": 6.95, "w_in": 10.3, "h_in": 0.35},
        name="disclaimer",
        placeholder_text="{{disclaimer}}",
        font=FONT_BODY, size_pt=10, bold=False,
        color=BRAND["dark_gray"], align="center"
    )
    return slide


def build_template_three_step_circles_dark(prs: Presentation,
                                            slide_data: dict, bg_image: str):
    """
    Slide 4: Three-step process with numbered circles
    - Dark charcoal background
    - Breadcrumb top-left (light text), icon top-right (white)
    - Large centered white title
    - Yellow italic subtitle
    - 3 numbered circle infographic (baked in background; labels editable)
    """
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    add_background_image(slide, bg_image)

    add_placeholder_textbox(slide,
        pos={"x_in": 0.35, "y_in": 0.18, "w_in": 3.0, "h_in": 0.32},
        name="breadcrumb",
        placeholder_text="{{breadcrumb}}",
        font=FONT_HEADER, size_pt=11, bold=False,
        all_caps=True, color=BRAND["white"], align="left"
    )
    add_placeholder_textbox(slide,
        pos={"x_in": 1.0, "y_in": 0.55, "w_in": 11.3, "h_in": 0.75},
        name="slide_title",
        placeholder_text="{{slide_title}}",
        font=FONT_HEADER, size_pt=40, bold=True,
        all_caps=True, color=BRAND["white"], align="center"
    )
    add_placeholder_textbox(slide,
        pos={"x_in": 1.0, "y_in": 1.35, "w_in": 11.3, "h_in": 0.45},
        name="subtitle",
        placeholder_text="{{subtitle}}",
        font=FONT_BODY, size_pt=16, bold=False,
        color=BRAND["yellow"], align="center"
    )
    # Step labels (3 circles — text only, circles are in background image)
    for i, label_name in enumerate(["step1_label", "step2_label", "step3_label"]):
        x = 0.7 + i * 4.2
        add_placeholder_textbox(slide,
            pos={"x_in": x, "y_in": 5.5, "w_in": 3.5, "h_in": 0.7},
            name=label_name,
            placeholder_text="{{" + label_name + "}}",
            font=FONT_BODY, size_pt=13, bold=True,
            color=BRAND["white"], align="center"
        )
    return slide


def build_template_three_col_equation_dark(prs: Presentation,
                                            slide_data: dict, bg_image: str):
    """
    Slide 8: Three-column equation (Base + Trigger = Composite)
    - Dark background
    - Multi-line title with yellow accent words
    - Three card columns with yellow pill headers
    """
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    add_background_image(slide, bg_image)

    add_placeholder_textbox(slide,
        pos={"x_in": 0.35, "y_in": 0.18, "w_in": 3.0, "h_in": 0.32},
        name="breadcrumb",
        placeholder_text="{{breadcrumb}}",
        font=FONT_HEADER, size_pt=11, bold=False,
        all_caps=True, color=BRAND["white"], align="left"
    )
    add_placeholder_textbox(slide,
        pos={"x_in": 0.5, "y_in": 0.5, "w_in": 12.3, "h_in": 1.2},
        name="slide_title",
        placeholder_text="{{slide_title}}",
        font=FONT_HEADER, size_pt=36, bold=True,
        all_caps=True, color=BRAND["white"], align="center"
    )
    col_positions = [
        {"x_in": 0.45, "y_in": 2.0, "w_in": 3.4, "h_in": 2.8},
        {"x_in": 4.75, "y_in": 2.0, "w_in": 3.4, "h_in": 2.8},
        {"x_in": 9.05, "y_in": 2.0, "w_in": 3.4, "h_in": 2.8},
    ]
    col_names = ["col1_header", "col2_header", "col3_header"]
    body_names = ["col1_body", "col2_body", "col3_body"]
    for i, (cp, cn, bn) in enumerate(zip(col_positions, col_names, body_names)):
        header_pos = {**cp, "y_in": cp["y_in"], "h_in": 0.38}
        body_pos   = {**cp, "y_in": cp["y_in"] + 0.55, "h_in": 2.1}
        add_placeholder_textbox(slide,
            pos=header_pos, name=cn,
            placeholder_text="{{" + cn + "}}",
            font=FONT_BODY, size_pt=14, bold=True,
            color=BRAND["dark_gray"], align="center"
        )
        add_placeholder_textbox(slide,
            pos=body_pos, name=bn,
            placeholder_text="{{" + bn + "}}",
            font=FONT_BODY, size_pt=13, bold=False,
            color=BRAND["white"], align="center"
        )
    return slide


def build_template_four_col_equation_dark(prs: Presentation,
                                           slide_data: dict, bg_image: str):
    """
    Slide 9: Four-column equation (Base + Trigger + Trigger = Composite)
    - Dark background
    - Title + yellow subtitle
    - Four card columns, last one yellow-highlighted
    """
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    add_background_image(slide, bg_image)

    add_placeholder_textbox(slide,
        pos={"x_in": 0.35, "y_in": 0.18, "w_in": 3.0, "h_in": 0.32},
        name="breadcrumb",
        placeholder_text="{{breadcrumb}}",
        font=FONT_HEADER, size_pt=11, bold=False,
        all_caps=True, color=BRAND["white"], align="left"
    )
    add_placeholder_textbox(slide,
        pos={"x_in": 0.5, "y_in": 0.5, "w_in": 12.3, "h_in": 0.65},
        name="slide_title",
        placeholder_text="{{slide_title}}",
        font=FONT_HEADER, size_pt=36, bold=True,
        all_caps=True, color=BRAND["white"], align="center"
    )
    add_placeholder_textbox(slide,
        pos={"x_in": 0.5, "y_in": 1.2, "w_in": 12.3, "h_in": 0.38},
        name="subtitle",
        placeholder_text="{{subtitle}}",
        font=FONT_BODY, size_pt=15, bold=False,
        color=BRAND["yellow"], align="center"
    )
    cols = [
        (0.25, "col1_header", "col1_body"),
        (3.35, "col2_header", "col2_body"),
        (6.45, "col3_header", "col3_body"),
        (9.55, "col4_header", "col4_body"),
    ]
    for x_start, ch, cb in cols:
        add_placeholder_textbox(slide,
            pos={"x_in": x_start, "y_in": 2.0, "w_in": 2.9, "h_in": 0.38},
            name=ch, placeholder_text="{{" + ch + "}}",
            font=FONT_BODY, size_pt=14, bold=True,
            color=BRAND["dark_gray"], align="center"
        )
        add_placeholder_textbox(slide,
            pos={"x_in": x_start, "y_in": 2.55, "w_in": 2.9, "h_in": 2.3},
            name=cb, placeholder_text="{{" + cb + "}}",
            font=FONT_BODY, size_pt=13, bold=False,
            color=BRAND["white"], align="center"
        )
    return slide


def build_template_data_table_light(prs: Presentation, slide_data: dict,
                                     bg_image: str):
    """
    Slide 11: Full-width data comparison table
    - Light background
    - Top-left title (not centered)
    - Full-width table with yellow header row
    - Source footnote
    """
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    add_background_image(slide, bg_image)

    add_placeholder_textbox(slide,
        pos={"x_in": 0.35, "y_in": 0.18, "w_in": 12.5, "h_in": 0.45},
        name="slide_title",
        placeholder_text="{{slide_title}}",
        font=FONT_HEADER, size_pt=16, bold=False,
        all_caps=True, color=BRAND["dark_gray"], align="left"
    )
    # Table area — content baked in background; add named zone for assembler
    add_placeholder_image_zone(slide,
        pos={"x_in": 0.35, "y_in": 0.8, "w_in": 12.6, "h_in": 5.9},
        name="table_area"
    )
    add_placeholder_textbox(slide,
        pos={"x_in": 0.35, "y_in": 6.95, "w_in": 12.6, "h_in": 0.35},
        name="footnote",
        placeholder_text="{{footnote}}",
        font=FONT_BODY, size_pt=10, bold=False,
        color=BRAND["dark_gray"], align="center"
    )
    return slide


def build_template_comparison_table(prs: Presentation, slide_data: dict,
                                     bg_image: str):
    """
    Slide 12: Two-row comparison table with yellow highlight
    - Light background
    - Breadcrumb + centered title
    - Small 2-row table with yellow-highlighted Potomac row
    """
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    add_background_image(slide, bg_image)

    add_placeholder_textbox(slide,
        pos={"x_in": 0.35, "y_in": 0.18, "w_in": 6.0, "h_in": 0.32},
        name="breadcrumb",
        placeholder_text="{{breadcrumb}}",
        font=FONT_HEADER, size_pt=11, bold=False,
        all_caps=True, color=BRAND["dark_gray"], align="left"
    )
    add_placeholder_textbox(slide,
        pos={"x_in": 1.0, "y_in": 0.55, "w_in": 11.3, "h_in": 0.65},
        name="slide_title",
        placeholder_text="{{slide_title}}",
        font=FONT_HEADER, size_pt=36, bold=True,
        all_caps=True, color=BRAND["dark_gray"], align="center"
    )
    add_placeholder_image_zone(slide,
        pos={"x_in": 0.8, "y_in": 1.8, "w_in": 11.7, "h_in": 3.2},
        name="table_area"
    )
    add_placeholder_textbox(slide,
        pos={"x_in": 0.5, "y_in": 6.95, "w_in": 12.3, "h_in": 0.35},
        name="footnote",
        placeholder_text="{{footnote}}",
        font=FONT_BODY, size_pt=10, bold=False,
        color=BRAND["dark_gray"], align="center"
    )
    return slide


def build_template_split_dark_light(prs: Presentation, slide_data: dict,
                                     bg_image: str):
    """
    Slide 10: Diagonal split — dark left, light right
    - Left side: heading + body text
    - Right side: allocation charts / infographic
    """
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    add_background_image(slide, bg_image)

    add_placeholder_textbox(slide,
        pos={"x_in": 0.35, "y_in": 0.18, "w_in": 4.0, "h_in": 0.32},
        name="breadcrumb",
        placeholder_text="{{breadcrumb}}",
        font=FONT_HEADER, size_pt=11, bold=False,
        all_caps=True, color=BRAND["white"], align="left"
    )
    add_placeholder_textbox(slide,
        pos={"x_in": 0.35, "y_in": 1.2, "w_in": 5.5, "h_in": 0.75},
        name="left_heading",
        placeholder_text="{{left_heading}}",
        font=FONT_HEADER, size_pt=40, bold=True,
        all_caps=True, color=BRAND["white"], align="left"
    )
    add_placeholder_textbox(slide,
        pos={"x_in": 0.35, "y_in": 2.1, "w_in": 5.5, "h_in": 3.2},
        name="left_body",
        placeholder_text="{{left_body}}",
        font=FONT_BODY, size_pt=14, bold=False,
        color=BRAND["white"], align="left"
    )
    add_placeholder_image_zone(slide,
        pos={"x_in": 6.5, "y_in": 0.5, "w_in": 6.5, "h_in": 7.0},
        name="right_infographic"
    )
    return slide


def build_template_three_circles_dark(prs: Presentation, slide_data: dict,
                                       bg_image: str):
    """
    Slide 13: Three circles on dark background (no numbers)
    - "How are advisors using X?"
    """
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    add_background_image(slide, bg_image)

    add_placeholder_textbox(slide,
        pos={"x_in": 0.35, "y_in": 0.18, "w_in": 6.0, "h_in": 0.32},
        name="breadcrumb",
        placeholder_text="{{breadcrumb}}",
        font=FONT_HEADER, size_pt=11, bold=False,
        all_caps=True, color=BRAND["white"], align="left"
    )
    add_placeholder_textbox(slide,
        pos={"x_in": 1.0, "y_in": 0.55, "w_in": 11.3, "h_in": 0.75},
        name="slide_title",
        placeholder_text="{{slide_title}}",
        font=FONT_HEADER, size_pt=40, bold=True,
        all_caps=True, color=BRAND["white"], align="center"
    )
    for i, name in enumerate(["circle1_text", "circle2_text", "circle3_text"]):
        x = 0.5 + i * 4.1
        add_placeholder_textbox(slide,
            pos={"x_in": x, "y_in": 2.2, "w_in": 3.7, "h_in": 2.5},
            name=name, placeholder_text="{{" + name + "}}",
            font=FONT_BODY, size_pt=15, bold=False,
            color=BRAND["white"], align="center"
        )
    return slide


def build_template_thank_you(prs: Presentation, slide_data: dict,
                              bg_image: str):
    """Slide 14: Thank you / contact slide."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    add_background_image(slide, bg_image)

    add_placeholder_textbox(slide,
        pos={"x_in": 0.35, "y_in": 0.18, "w_in": 3.0, "h_in": 0.32},
        name="breadcrumb",
        placeholder_text="{{breadcrumb}}",
        font=FONT_HEADER, size_pt=11, bold=False,
        all_caps=True, color=BRAND["white"], align="left"
    )
    add_placeholder_textbox(slide,
        pos={"x_in": 1.0, "y_in": 0.6, "w_in": 11.3, "h_in": 0.75},
        name="heading",
        placeholder_text="{{heading}}",
        font=FONT_HEADER, size_pt=44, bold=True,
        all_caps=True, color=BRAND["white"], align="center"
    )
    add_placeholder_textbox(slide,
        pos={"x_in": 1.0, "y_in": 1.45, "w_in": 11.3, "h_in": 0.45},
        name="subheading",
        placeholder_text="{{subheading}}",
        font=FONT_BODY, size_pt=18, bold=False,
        color=BRAND["yellow"], align="center"
    )
    add_placeholder_image_zone(slide,
        pos={"x_in": 0.5, "y_in": 2.1, "w_in": 6.8, "h_in": 4.8},
        name="map_image"
    )
    add_placeholder_textbox(slide,
        pos={"x_in": 7.8, "y_in": 2.1, "w_in": 5.0, "h_in": 0.4},
        name="contact_url",
        placeholder_text="{{contact_url}}",
        font=FONT_BODY, size_pt=16, bold=False,
        color=BRAND["white"], align="left"
    )
    add_placeholder_image_zone(slide,
        pos={"x_in": 7.8, "y_in": 2.7, "w_in": 2.5, "h_in": 2.5},
        name="qr_code"
    )
    return slide


def build_template_text_dark(prs: Presentation, slide_data: dict,
                              bg_image: str):
    """
    Slides 15,16: Dense text slides (Disclosures / Definitions)
    - Dark background
    - Top-left title
    - Large body text area
    """
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    add_background_image(slide, bg_image)

    add_placeholder_textbox(slide,
        pos={"x_in": 0.35, "y_in": 0.18, "w_in": 12.0, "h_in": 0.42},
        name="slide_title",
        placeholder_text="{{slide_title}}",
        font=FONT_HEADER, size_pt=18, bold=False,
        all_caps=True, color=BRAND["white"], align="left"
    )
    add_placeholder_textbox(slide,
        pos={"x_in": 0.35, "y_in": 0.75, "w_in": 12.6, "h_in": 6.4},
        name="body_text",
        placeholder_text="{{body_text}}",
        font=FONT_BODY, size_pt=11, bold=False,
        color=BRAND["white"], align="left"
    )
    return slide


def build_template_data_flow_diagram_light(prs: Presentation,
                                            slide_data: dict, bg_image: str):
    """
    Slide 3: Data flow diagram (5-box, center yellow)
    - Light background
    - Large centered title
    - Diagram is baked into background; only labels are editable
    """
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    add_background_image(slide, bg_image)

    add_placeholder_textbox(slide,
        pos={"x_in": 0.35, "y_in": 0.18, "w_in": 3.0, "h_in": 0.32},
        name="breadcrumb",
        placeholder_text="{{breadcrumb}}",
        font=FONT_HEADER, size_pt=11, bold=False,
        all_caps=True, color=BRAND["dark_gray"], align="left"
    )
    add_placeholder_textbox(slide,
        pos={"x_in": 0.5, "y_in": 0.55, "w_in": 12.3, "h_in": 0.65},
        name="slide_title",
        placeholder_text="{{slide_title}}",
        font=FONT_HEADER, size_pt=36, bold=True,
        all_caps=True, color=BRAND["dark_gray"], align="center"
    )
    box_defs = [
        ("top_left_header",  "top_left_body",  0.25, 1.55),
        ("top_right_header", "top_right_body",  9.3, 1.55),
        ("center_header",    "center_body",     4.7, 1.55),
        ("btm_left_header",  "btm_left_body",   0.25, 4.45),
        ("btm_right_header", "btm_right_body",  9.3, 4.45),
    ]
    for hname, bname, x, y in box_defs:
        add_placeholder_textbox(slide,
            pos={"x_in": x, "y_in": y, "w_in": 3.5, "h_in": 0.38},
            name=hname, placeholder_text="{{" + hname + "}}",
            font=FONT_BODY, size_pt=14, bold=True,
            color=BRAND["dark_gray"], align="center"
        )
        add_placeholder_textbox(slide,
            pos={"x_in": x, "y_in": y + 0.5, "w_in": 3.5, "h_in": 2.5},
            name=bname, placeholder_text="{{" + bname + "}}",
            font=FONT_BODY, size_pt=12, bold=False,
            color=BRAND["white"], align="center"
        )
    return slide


# ─── TEMPLATE REGISTRY ────────────────────────────────────────────────────────

TEMPLATE_BUILDERS = {
    "cover-hero":                  build_template_cover_hero,
    "strategy-intro":              build_template_strategy_intro,
    "chart-full-light":            build_template_chart_full_light,
    "three-step-circles-dark":     build_template_three_step_circles_dark,
    "three-col-equation-dark":     build_template_three_col_equation_dark,
    "four-col-equation-dark":      build_template_four_col_equation_dark,
    "data-table-light":            build_template_data_table_light,
    "comparison-table":            build_template_comparison_table,
    "split-dark-light":            build_template_split_dark_light,
    "three-circles-dark":          build_template_three_circles_dark,
    "thank-you":                   build_template_thank_you,
    "text-dark":                   build_template_text_dark,
    "data-flow-diagram-light":     build_template_data_flow_diagram_light,
}

# Map from slide number in Bull Bear to template_id
# You can override these in the manifest JSON under each slide's "template_id"
BULL_BEAR_SLIDE_MAP = {
    1:  "cover-hero",
    2:  "strategy-intro",
    3:  "data-flow-diagram-light",
    4:  "three-step-circles-dark",
    5:  "chart-full-light",
    6:  "chart-full-light",
    7:  "chart-full-light",
    8:  "three-col-equation-dark",
    9:  "four-col-equation-dark",
    10: "split-dark-light",
    11: "data-table-light",
    12: "comparison-table",
    13: "three-circles-dark",
    14: "thank-you",
    15: "text-dark",
    16: "text-dark",
}


# ─── MAIN ────────────────────────────────────────────────────────────────────

def convert(manifest_path: str, images_dir: str, output_dir: str,
            logos_dir: str = None):
    output_path = Path(output_dir)
    output_path.mkdir(parents=True, exist_ok=True)

    with open(manifest_path, "r", encoding="utf-8") as f:
        manifest = json.loads(f.read(), strict=False)

    print(f"Loaded manifest: {manifest['document_name']}")
    print(f"Slides: {manifest['total_slides']}")
    print(f"Dimensions: {manifest['slide_width_in']}\" x {manifest['slide_height_in']}\"")

    slide_w = Inches(manifest["slide_width_in"])
    slide_h = Inches(manifest["slide_height_in"])

    # Update global dimensions if different from default
    global SLIDE_W, SLIDE_H
    SLIDE_W = slide_w
    SLIDE_H = slide_h

    # Build one .pptx per template type, deduplicated
    built_templates = {}

    for slide_data in manifest["slides"]:
        slide_num  = slide_data["slide_number"]
        template_id = slide_data.get("template_id") or \
                      BULL_BEAR_SLIDE_MAP.get(slide_num)

        if not template_id:
            print(f"  [WARN] Slide {slide_num}: no template_id, skipping")
            continue

        if template_id in built_templates:
            print(f"  [OK] Slide {slide_num}: template '{template_id}' already built")
            continue

        builder = TEMPLATE_BUILDERS.get(template_id)
        if not builder:
            print(f"  [WARN] Slide {slide_num}: unknown template_id '{template_id}'")
            continue

        # Find reference image
        bg_image = os.path.join(images_dir,
                                f"slide_{str(slide_num).zfill(2)}.png")
        if not os.path.exists(bg_image):
            print(f"  [WARN] Background image not found: {bg_image}")
            bg_image = None

        # Build presentation with single slide
        prs = Presentation()
        prs.slide_width  = slide_w
        prs.slide_height = slide_h

        print(f"  -> Building template '{template_id}' from slide {slide_num}...")
        try:
            builder(prs, slide_data, bg_image)
        except Exception as e:
            print(f"     ERROR: {e}")
            continue

        out_file = output_path / f"{template_id}.pptx"
        prs.save(str(out_file))
        built_templates[template_id] = str(out_file)
        print(f"     Saved: {out_file}")

    # Write summary manifest
    summary = {
        "deck_family":   manifest["deck_family"],
        "slide_w_in":    manifest["slide_width_in"],
        "slide_h_in":    manifest["slide_height_in"],
        "templates":     built_templates,
        "slide_map":     BULL_BEAR_SLIDE_MAP,
    }
    summary_path = output_path / "template_registry.json"
    with open(str(summary_path), "w") as f:
        json.dump(summary, f, indent=2)
    print(f"\nDone! {len(built_templates)} templates built.")
    print(f"   Registry: {summary_path}")


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Convert InDesign manifest to PPTX templates")
    parser.add_argument("--manifest", required=True)
    parser.add_argument("--images",   required=True)
    parser.add_argument("--output",   required=True)
    parser.add_argument("--logos",    default=None)
    args = parser.parse_args()

    convert(args.manifest, args.images, args.output, args.logos)
