"""
Potomac InDesign MCP Server
===========================
A local MCP server that controls Adobe InDesign via Windows COM automation.
Cline connects to this server and can trigger InDesign exports, run scripts,
and then chain into the Python template converter — all without manual steps.

REQUIREMENTS (run once):
    pip install mcp pywin32 Pillow python-pptx

REGISTER IN CLINE:
    Add to Cline MCP settings (see cline-mcp-config.json in this folder).

HOW IT WORKS:
    - Runs as a local stdio MCP server on your Windows machine
    - Uses win32com.client to open InDesign and dispatch COM commands
    - Executes the JSX manifest-export script inside InDesign
    - Runs the Python template converter after export completes
    - Cline orchestrates it all with natural language
"""

import asyncio
import json
import os
import sys
import subprocess
import tempfile
import time
from pathlib import Path
from typing import Optional

import mcp.server.stdio
import mcp.types as types
from mcp.server import Server, NotificationOptions
from mcp.server.models import InitializationOptions

# ─── SERVER SETUP ────────────────────────────────────────────────────────────

server = Server("indesign-mcp")

# Path to this file's directory (where converter scripts live alongside it)
SERVER_DIR = Path(__file__).parent.resolve()
CONVERTER_SCRIPT = SERVER_DIR / "convert_to_templates.py"
JSX_SCRIPT       = SERVER_DIR / "export_manifest.jsx"


# ─── HELPERS ─────────────────────────────────────────────────────────────────

def get_indesign_com():
    """Get a COM handle to the running (or newly launched) InDesign instance."""
    errors = []
    try:
        import win32com.client
        import pythoncom
        pythoncom.CoInitialize()

        # Try to connect to a running instance first
        try:
            app = win32com.client.GetActiveObject("InDesign.Application")
            return app, None
        except Exception as e:
            errors.append(f"GetActiveObject failed: {e}")

        # Launch InDesign — try common version strings (newest first)
        version_strings = [
            "InDesign.Application.2026",
            "InDesign.Application.2025",
            "InDesign.Application.2024",
            "InDesign.Application.2023",
            "InDesign.Application.2022",
            "InDesign.Application",
        ]
        for vs in version_strings:
            try:
                app = win32com.client.Dispatch(vs)
                try:
                    app.Visible = True
                except Exception:
                    pass  # Visible may fail but COM handle is valid
                time.sleep(3)  # give InDesign time to load
                return app, None
            except Exception as e:
                errors.append(f"Dispatch('{vs}') failed: {e}")
                continue

        return None, f"Could not launch InDesign. Tried: {'; '.join(errors)}"

    except ImportError as e:
        return None, f"pywin32 not installed ({e}). Run: pip install pywin32"
    except Exception as e:
        return None, f"COM error: {e}. Previous errors: {'; '.join(errors)}"


def _get_undo_mode(app):
    """Get the idAutoUndo enum value compatible with the current InDesign version."""
    try:
        # Try using EnsureDispatch for typed constants
        import win32com.client
        typed_app = win32com.client.gencache.EnsureDispatch(app)
        return win32com.client.constants.idAutoUndo
    except Exception:
        pass
    # Fallback: the raw enum value for idAutoUndo
    # 1699116885 = 0x65557544 (idAutoUndo in newer InDesign versions)
    return 1699116885


def run_jsx_in_indesign(app, jsx_path: str) -> tuple[bool, str]:
    """Execute a JSX script inside InDesign via COM."""
    try:
        import win32com.client

        with open(jsx_path, "r", encoding="utf-8") as f:
            script_content = f.read()

        # ScriptLanguage.idJavascript = 1246973031
        # Try without UndoMode first (safest for InDesign 2026+)
        try:
            result = app.DoScript(script_content, 1246973031)
            return True, str(result) if result else "Script completed"
        except Exception as e1:
            # Fallback: try with empty args
            try:
                result = app.DoScript(script_content, 1246973031, [])
                return True, str(result) if result else "Script completed"
            except Exception as e2:
                return False, f"Script error: attempt1={e1}, attempt2={e2}"

    except Exception as e:
        return False, f"Script error: {e}"


def run_jsx_string_in_indesign(app, jsx_code: str) -> tuple[bool, str]:
    """Execute a JSX code string inside InDesign via COM."""
    try:
        # Try without UndoMode (safest for InDesign 2026+)
        try:
            result = app.DoScript(jsx_code, 1246973031)
            return True, str(result) if result else "OK"
        except Exception:
            result = app.DoScript(jsx_code, 1246973031, [])
            return True, str(result) if result else "OK"
    except Exception as e:
        return False, f"Error: {e}"


# ─── TOOL DEFINITIONS ────────────────────────────────────────────────────────

@server.list_tools()
async def handle_list_tools() -> list[types.Tool]:
    return [
        types.Tool(
            name="indesign_open_document",
            description=(
                "Opens an InDesign (.indd) file. Call this before running "
                "export scripts. Returns the document name and page count."
            ),
            inputSchema={
                "type": "object",
                "properties": {
                    "indd_path": {
                        "type": "string",
                        "description": "Full Windows path to the .indd file, e.g. C:\\Users\\...\\Bull_Bear.indd"
                    }
                },
                "required": ["indd_path"]
            }
        ),
        types.Tool(
            name="indesign_export_manifest",
            description=(
                "Runs the Potomac manifest export script on the currently open "
                "InDesign document. Extracts all text frames, shapes, and image "
                "zones from every page, exports high-res PNGs, and writes a "
                "JSON manifest. This is the InDesign → JSON step of the "
                "template extraction pipeline. Takes 3-8 minutes for large decks."
            ),
            inputSchema={
                "type": "object",
                "properties": {
                    "indd_path": {
                        "type": "string",
                        "description": "Full path to the .indd file to export"
                    },
                    "output_dir": {
                        "type": "string",
                        "description": (
                            "Directory where manifest JSON and slide PNGs will be saved. "
                            "Defaults to a 'potomac_export' folder next to the .indd file."
                        )
                    }
                },
                "required": ["indd_path"]
            }
        ),
        types.Tool(
            name="indesign_run_jsx",
            description=(
                "Runs a custom JSX ExtendScript inside InDesign. "
                "Use for one-off queries like getting document info, "
                "checking font names, or page dimensions."
            ),
            inputSchema={
                "type": "object",
                "properties": {
                    "jsx_code": {
                        "type": "string",
                        "description": "ExtendScript JSX code to execute"
                    }
                },
                "required": ["jsx_code"]
            }
        ),
        types.Tool(
            name="indesign_get_document_info",
            description=(
                "Returns metadata about the currently open InDesign document: "
                "name, page count, dimensions, fonts used, layers, and color swatches."
            ),
            inputSchema={
                "type": "object",
                "properties": {},
                "required": []
            }
        ),
        types.Tool(
            name="indesign_list_open_documents",
            description="Lists all currently open InDesign documents.",
            inputSchema={
                "type": "object",
                "properties": {},
                "required": []
            }
        ),
        types.Tool(
            name="convert_manifest_to_templates",
            description=(
                "Runs the Python converter that reads the InDesign manifest JSON "
                "and builds .pptx template files. This is Step 2 of the pipeline "
                "(after indesign_export_manifest). Produces one .pptx per slide "
                "type in the output directory."
            ),
            inputSchema={
                "type": "object",
                "properties": {
                    "manifest_path": {
                        "type": "string",
                        "description": "Full path to the *_manifest.json file from the InDesign export"
                    },
                    "images_dir": {
                        "type": "string",
                        "description": "Path to the slide_images/ folder from the InDesign export"
                    },
                    "output_dir": {
                        "type": "string",
                        "description": "Where to save the .pptx template files. Should be inside your backend's pptx_engine/templates/{deck-family}/ folder."
                    }
                },
                "required": ["manifest_path", "images_dir", "output_dir"]
            }
        ),
        types.Tool(
            name="run_full_pipeline",
            description=(
                "ONE-SHOT: Runs the entire InDesign → PPTX template pipeline. "
                "Opens the .indd file, runs the manifest export, waits for "
                "completion, then runs the Python converter. "
                "Returns paths to all generated template files."
            ),
            inputSchema={
                "type": "object",
                "properties": {
                    "indd_path": {
                        "type": "string",
                        "description": "Full path to the .indd file"
                    },
                    "backend_templates_dir": {
                        "type": "string",
                        "description": (
                            "Path to the backend's templates folder, e.g. "
                            "C:\\Users\\SohaibAli\\PycharmProjects\\Potomac-Analyst-Workbench\\pptx_engine\\templates"
                        )
                    },
                    "deck_family": {
                        "type": "string",
                        "description": "Slug for this deck, e.g. 'bull-bear', 'guardian'"
                    }
                },
                "required": ["indd_path", "backend_templates_dir", "deck_family"]
            }
        ),
        types.Tool(
            name="verify_templates",
            description=(
                "Checks that all expected .pptx template files were created "
                "correctly. Opens each one, counts slides, and lists the named "
                "placeholder shapes found. Use after convert_manifest_to_templates "
                "to confirm the output is good."
            ),
            inputSchema={
                "type": "object",
                "properties": {
                    "templates_dir": {
                        "type": "string",
                        "description": "Path to the deck-family templates folder"
                    }
                },
                "required": ["templates_dir"]
            }
        ),
    ]


# ─── TOOL HANDLERS ───────────────────────────────────────────────────────────

@server.call_tool()
async def handle_call_tool(name: str, arguments: dict) -> list[types.TextContent]:

    # ── indesign_open_document ───────────────────────────────────────────────
    if name == "indesign_open_document":
        indd_path = arguments["indd_path"]

        if not Path(indd_path).exists():
            return [types.TextContent(type="text",
                text=f"ERROR: File not found: {indd_path}")]

        app, err = get_indesign_com()
        if err:
            return [types.TextContent(type="text", text=f"ERROR: {err}")]

        try:
            # Check if already open
            for i in range(app.Documents.Count):
                doc = app.Documents.Item(i + 1)
                if doc.FullName.lower() == indd_path.lower():
                    return [types.TextContent(type="text",
                        text=json.dumps({
                            "status": "already_open",
                            "name": doc.Name,
                            "pages": doc.Pages.Count
                        }, indent=2))]

            # Open the document
            doc = app.Open(indd_path)
            time.sleep(1)

            return [types.TextContent(type="text",
                text=json.dumps({
                    "status": "opened",
                    "name": doc.Name,
                    "full_path": doc.FullName,
                    "pages": doc.Pages.Count,
                }, indent=2))]

        except Exception as e:
            return [types.TextContent(type="text", text=f"ERROR opening document: {e}")]

    # ── indesign_get_document_info ───────────────────────────────────────────
    elif name == "indesign_get_document_info":
        app, err = get_indesign_com()
        if err:
            return [types.TextContent(type="text", text=f"ERROR: {err}")]

        jsx = """
(function() {
    if (app.documents.length === 0) return JSON.stringify({error: "No document open"});
    var doc = app.activeDocument;
    var page = doc.pages[0];
    var w = page.bounds[3] - page.bounds[1];
    var h = page.bounds[2] - page.bounds[0];
    var fonts = [];
    for (var i = 0; i < doc.fonts.length && i < 20; i++) {
        fonts.push(doc.fonts[i].name);
    }
    var layers = [];
    for (var i = 0; i < doc.layers.length; i++) {
        layers.push(doc.layers[i].name);
    }
    var swatches = [];
    for (var i = 0; i < doc.swatches.length && i < 30; i++) {
        swatches.push(doc.swatches[i].name);
    }
    return JSON.stringify({
        name: doc.name,
        fullPath: doc.filePath.absoluteURI,
        pageCount: doc.pages.length,
        widthPt: w,
        heightPt: h,
        widthIn: w/72,
        heightIn: h/72,
        fonts: fonts,
        layers: layers,
        swatches: swatches
    });
})();
"""
        ok, result = run_jsx_string_in_indesign(app, jsx)
        return [types.TextContent(type="text",
            text=result if ok else f"ERROR: {result}")]

    # ── indesign_list_open_documents ─────────────────────────────────────────
    elif name == "indesign_list_open_documents":
        app, err = get_indesign_com()
        if err:
            return [types.TextContent(type="text", text=f"ERROR: {err}")]

        try:
            docs = []
            for i in range(app.Documents.Count):
                doc = app.Documents.Item(i + 1)
                docs.append({
                    "name": doc.Name,
                    "path": doc.FullName,
                    "pages": doc.Pages.Count,
                    "modified": doc.Modified
                })
            return [types.TextContent(type="text",
                text=json.dumps({"open_documents": docs}, indent=2))]
        except Exception as e:
            return [types.TextContent(type="text", text=f"ERROR: {e}")]

    # ── indesign_run_jsx ─────────────────────────────────────────────────────
    elif name == "indesign_run_jsx":
        app, err = get_indesign_com()
        if err:
            return [types.TextContent(type="text", text=f"ERROR: {err}")]

        ok, result = run_jsx_string_in_indesign(app, arguments["jsx_code"])
        return [types.TextContent(type="text",
            text=result if ok else f"ERROR: {result}")]

    # ── indesign_export_manifest ─────────────────────────────────────────────
    elif name == "indesign_export_manifest":
        indd_path  = arguments["indd_path"]
        indd_path_obj = Path(indd_path)

        if not indd_path_obj.exists():
            return [types.TextContent(type="text",
                text=f"ERROR: File not found: {indd_path}")]

        output_dir = arguments.get("output_dir") or \
                     str(indd_path_obj.parent / "potomac_export")
        Path(output_dir).mkdir(parents=True, exist_ok=True)

        # First open the document
        app, err = get_indesign_com()
        if err:
            return [types.TextContent(type="text", text=f"ERROR: {err}")]

        # Open if not already open
        try:
            doc_open = False
            for i in range(app.Documents.Count):
                d = app.Documents.Item(i + 1)
                if d.FullName.lower() == str(indd_path).lower():
                    doc_open = True
                    app.ActiveDocument = d
                    break
            if not doc_open:
                app.Open(indd_path)
                time.sleep(2)
        except Exception as e:
            return [types.TextContent(type="text",
                text=f"ERROR opening document: {e}")]

        if not JSX_SCRIPT.exists():
            return [types.TextContent(type="text",
                text=f"ERROR: export_manifest.jsx not found at {JSX_SCRIPT}. "
                     f"Place it in the same folder as this MCP server.")]

        # Read the JSX script and inject the output_dir
        with open(JSX_SCRIPT, "r", encoding="utf-8") as f:
            jsx_content = f.read()

        # Patch the output dir into the script
        output_dir_escaped = output_dir.replace("\\", "/")
        patched_jsx = f"""
// Injected by MCP server
var INJECTED_OUTPUT_DIR = "{output_dir_escaped}";
""" + jsx_content.replace(
    'var outputFolder = new Folder(docPath + "/potomac_export");',
    'var outputFolder = new Folder(INJECTED_OUTPUT_DIR || (docPath + "/potomac_export"));'
)

        # Write patched JSX to temp file
        tmp_jsx = Path(tempfile.gettempdir()) / "potomac_export_patched.jsx"
        with open(tmp_jsx, "w", encoding="utf-8") as f:
            f.write(patched_jsx)

        # Run the script (this will take several minutes)
        ok, result = run_jsx_in_indesign(app, str(tmp_jsx))

        if not ok:
            return [types.TextContent(type="text", text=f"ERROR running script: {result}")]

        # Find the manifest
        doc_name = Path(indd_path).stem
        manifest_path = Path(output_dir) / f"{doc_name}_manifest.json"

        if manifest_path.exists():
            with open(manifest_path) as f:
                manifest_summary = json.load(f)

            return [types.TextContent(type="text",
                text=json.dumps({
                    "status": "success",
                    "manifest_path": str(manifest_path),
                    "images_dir": str(Path(output_dir) / "slide_images"),
                    "slide_count": manifest_summary.get("total_slides", "?"),
                    "document_name": manifest_summary.get("document_name", ""),
                    "next_step": (
                        f"Now call convert_manifest_to_templates with "
                        f"manifest_path='{manifest_path}' and "
                        f"images_dir='{Path(output_dir) / 'slide_images'}'"
                    )
                }, indent=2))]
        else:
            return [types.TextContent(type="text",
                text=json.dumps({
                    "status": "script_ran_but_manifest_missing",
                    "output_dir": output_dir,
                    "script_result": result,
                    "note": "Check InDesign for error dialogs."
                }, indent=2))]

    # ── convert_manifest_to_templates ────────────────────────────────────────
    elif name == "convert_manifest_to_templates":
        manifest_path = arguments["manifest_path"]
        images_dir    = arguments["images_dir"]
        output_dir    = arguments["output_dir"]

        if not Path(manifest_path).exists():
            return [types.TextContent(type="text",
                text=f"ERROR: Manifest not found: {manifest_path}")]

        if not CONVERTER_SCRIPT.exists():
            return [types.TextContent(type="text",
                text=f"ERROR: convert_to_templates.py not found at {CONVERTER_SCRIPT}")]

        Path(output_dir).mkdir(parents=True, exist_ok=True)

        result = subprocess.run(
            [
                sys.executable, str(CONVERTER_SCRIPT),
                "--manifest", manifest_path,
                "--images",   images_dir,
                "--output",   output_dir,
            ],
            capture_output=True,
            text=True,
            timeout=300  # 5 min max
        )

        # List what was created
        created = list(Path(output_dir).glob("*.pptx"))

        return [types.TextContent(type="text",
            text=json.dumps({
                "status": "success" if result.returncode == 0 else "error",
                "returncode": result.returncode,
                "stdout": result.stdout[-3000:] if result.stdout else "",
                "stderr": result.stderr[-1000:] if result.stderr else "",
                "templates_created": [f.name for f in created],
                "templates_dir": output_dir,
                "count": len(created)
            }, indent=2))]

    # ── run_full_pipeline ────────────────────────────────────────────────────
    elif name == "run_full_pipeline":
        indd_path            = arguments["indd_path"]
        backend_templates_dir = arguments["backend_templates_dir"]
        deck_family          = arguments["deck_family"]

        indd = Path(indd_path)
        if not indd.exists():
            return [types.TextContent(type="text",
                text=f"ERROR: File not found: {indd_path}")]

        export_dir   = indd.parent / "potomac_export"
        manifest_path = export_dir / f"{indd.stem}_manifest.json"
        images_dir   = export_dir / "slide_images"
        output_dir   = Path(backend_templates_dir) / deck_family

        steps = []

        # Step 1: Open document
        app, err = get_indesign_com()
        if err:
            return [types.TextContent(type="text", text=f"ERROR: {err}")]

        try:
            doc_open = any(
                app.Documents.Item(i + 1).FullName.lower() == indd_path.lower()
                for i in range(app.Documents.Count)
            )
            if not doc_open:
                app.Open(indd_path)
                time.sleep(3)
            steps.append("✓ InDesign document opened")
        except Exception as e:
            return [types.TextContent(type="text",
                text=f"ERROR opening document: {e}")]

        # Step 2: Run export script
        steps.append("→ Running manifest export (this takes 3-8 minutes)...")

        export_dir.mkdir(parents=True, exist_ok=True)

        if not JSX_SCRIPT.exists():
            return [types.TextContent(type="text",
                text=f"ERROR: export_manifest.jsx not found at {JSX_SCRIPT}")]

        ok, result = run_jsx_in_indesign(app, str(JSX_SCRIPT))
        if not ok:
            return [types.TextContent(type="text",
                text=f"ERROR in export script: {result}")]
        steps.append(f"✓ Manifest exported to {export_dir}")

        # Step 3: Run converter
        if not manifest_path.exists():
            return [types.TextContent(type="text",
                text=f"ERROR: Manifest not created at expected path: {manifest_path}")]

        output_dir.mkdir(parents=True, exist_ok=True)
        steps.append(f"→ Running PPTX converter...")

        conv = subprocess.run(
            [sys.executable, str(CONVERTER_SCRIPT),
             "--manifest", str(manifest_path),
             "--images",   str(images_dir),
             "--output",   str(output_dir)],
            capture_output=True, text=True, timeout=300
        )

        created = list(output_dir.glob("*.pptx"))
        steps.append(f"✓ Created {len(created)} template files in {output_dir}")

        return [types.TextContent(type="text",
            text=json.dumps({
                "status": "complete",
                "steps": steps,
                "manifest_path": str(manifest_path),
                "templates_dir": str(output_dir),
                "templates_created": [f.name for f in created],
                "next_step": (
                    "Run verify_templates to confirm the output, "
                    "then commit the templates folder to git and push to Railway."
                )
            }, indent=2))]

    # ── verify_templates ─────────────────────────────────────────────────────
    elif name == "verify_templates":
        templates_dir = arguments["templates_dir"]
        pptx_files    = list(Path(templates_dir).glob("*.pptx"))

        if not pptx_files:
            return [types.TextContent(type="text",
                text=f"No .pptx files found in {templates_dir}")]

        try:
            from pptx import Presentation
        except ImportError:
            return [types.TextContent(type="text",
                text="ERROR: python-pptx not installed. Run: pip install python-pptx")]

        results = []
        for pptx_file in sorted(pptx_files):
            try:
                prs = Presentation(str(pptx_file))
                slide = prs.slides[0]
                placeholders = [s.name for s in slide.shapes
                                if s.name and not s.name.startswith("Picture")]
                results.append({
                    "file": pptx_file.name,
                    "slide_count": len(prs.slides),
                    "slide_width_in": round(prs.slide_width / 914400, 3),
                    "slide_height_in": round(prs.slide_height / 914400, 3),
                    "named_placeholders": placeholders,
                    "status": "✓ OK" if placeholders else "⚠ No named placeholders found"
                })
            except Exception as e:
                results.append({"file": pptx_file.name, "status": f"ERROR: {e}"})

        return [types.TextContent(type="text",
            text=json.dumps({
                "verified": len(results),
                "templates": results
            }, indent=2))]

    return [types.TextContent(type="text", text=f"Unknown tool: {name}")]


# ─── MAIN ────────────────────────────────────────────────────────────────────

async def main():
    async with mcp.server.stdio.stdio_server() as (read_stream, write_stream):
        await server.run(
            read_stream,
            write_stream,
            InitializationOptions(
                server_name="indesign-mcp",
                server_version="1.0.0",
                capabilities=server.get_capabilities(
                    notification_options=NotificationOptions(),
                    experimental_capabilities={}
                )
            )
        )

if __name__ == "__main__":
    asyncio.run(main())
