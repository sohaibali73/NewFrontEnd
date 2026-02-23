"""
Potomac PPTX Test Generator
============================
Enter a topic and this generates a Bull Bear presentation using
the pptxgenjs JSON templates extracted from InDesign.

Usage:
    python test_generate.py "Q1 2025 Market Outlook"
    python test_generate.py   (interactive prompt)

Output: test_output/<topic>.pptx
"""

import json
import os
import sys
import argparse
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    print("pip install python-pptx")
    sys.exit(1)

TEMPLATES_DIR = Path(__file__).parent.parent / "pptx_engine" / "templates" / "bull-bear"
OUTPUT_DIR = Path(__file__).parent / "test_output"

ALIGN_MAP = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}


def hex_to_rgb(h):
    if not h or len(h) < 6:
        return None
    h = h.lstrip("#")
    try:
        return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
    except:
        return None


def render_text_object(slide, obj_data):
    """Render a pptxgenjs text object to python-pptx."""
    text_data = obj_data.get("text", {})
    opts = text_data.get("options", {})

    x = Inches(opts.get("x", 0))
    y = Inches(opts.get("y", 0))
    w = Inches(opts.get("w", 1))
    h = Inches(opts.get("h", 0.5))

    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    txBox.fill.background()

    text_content = text_data.get("text", "")

    # Handle text array (multi-paragraph)
    if isinstance(text_content, list):
        for i, part in enumerate(text_content):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            part_opts = part.get("options", {})
            run = p.add_run()
            run.text = part.get("text", "")
            run.font.name = part_opts.get("fontFace", opts.get("fontFace", "Quicksand"))
            run.font.size = Pt(part_opts.get("fontSize", opts.get("fontSize", 12)))
            run.font.bold = part_opts.get("bold", False)
            run.font.italic = part_opts.get("italic", False)
            color = hex_to_rgb(part_opts.get("color", opts.get("color")))
            if color:
                run.font.color.rgb = color
            p.alignment = ALIGN_MAP.get(part_opts.get("align", opts.get("align", "left")), PP_ALIGN.LEFT)
    else:
        # Simple string
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = str(text_content)
        run.font.name = opts.get("fontFace", "Quicksand")
        run.font.size = Pt(opts.get("fontSize", 12))
        run.font.bold = opts.get("bold", False)
        color = hex_to_rgb(opts.get("color"))
        if color:
            run.font.color.rgb = color
        p.alignment = ALIGN_MAP.get(opts.get("align", "left"), PP_ALIGN.LEFT)


def render_rect_object(slide, obj_data):
    """Render a pptxgenjs rect object to python-pptx."""
    opts = obj_data.get("rect", {}).get("options", {})

    x = Inches(opts.get("x", 0))
    y = Inches(opts.get("y", 0))
    w = Inches(opts.get("w", 1))
    h = Inches(opts.get("h", 0.5))

    radius = opts.get("rectRadius", 0)
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius > 0 else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, x, y, w, h)

    fill = opts.get("fill", {})
    fill_color = hex_to_rgb(fill.get("color"))
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()

    line = opts.get("line", {})
    line_color = hex_to_rgb(line.get("color"))
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line.get("width", 1))
    else:
        shape.line.fill.background()


def render_image_object(slide, obj_data):
    """Render a pptxgenjs image object to python-pptx."""
    img_data = obj_data.get("image", {})
    opts = img_data.get("options", {})
    path = img_data.get("path", "")

    x = Inches(opts.get("x", 0))
    y = Inches(opts.get("y", 0))
    w = Inches(opts.get("w", 1))
    h = Inches(opts.get("h", 0.5))

    if path and os.path.exists(path):
        try:
            slide.shapes.add_picture(path, x, y, w, h)
            return
        except:
            pass

    # Placeholder rectangle for missing images
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.background()
    shape.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    shape.line.width = Pt(0.5)


def render_slide(prs, slide_template):
    """Render a single slide from pptxgenjs JSON template."""
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)

    # Background
    bg = slide_template.get("background", {})
    bg_color = hex_to_rgb(bg.get("color"))
    if bg_color:
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = bg_color

    # Render objects
    for obj in slide_template.get("objects", []):
        try:
            if "text" in obj:
                render_text_object(slide, obj)
            elif "rect" in obj:
                render_rect_object(slide, obj)
            elif "image" in obj:
                render_image_object(slide, obj)
        except Exception as e:
            pass  # skip failed objects

    return slide


def generate_presentation(topic: str, output_path: str = None):
    """Generate a Bull Bear presentation for the given topic."""

    # Load template catalog
    catalog_path = TEMPLATES_DIR / "template_catalog.json"
    if not catalog_path.exists():
        print(f"ERROR: Template catalog not found at {catalog_path}")
        print("Run the InDesign export pipeline first.")
        return None

    with open(catalog_path, "r", encoding="utf-8") as f:
        catalog = json.load(f)

    slide_w = catalog.get("slide_width_in", 13.333)
    slide_h = catalog.get("slide_height_in", 7.5)
    slides = catalog.get("slides", [])

    if not slides:
        print("ERROR: No slides in template catalog")
        return None

    print(f"Generating '{topic}' presentation...")
    print(f"  Template: {catalog['document_name']} ({len(slides)} slides)")
    print(f"  Size: {slide_w}\" x {slide_h}\"")

    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(slide_w)
    prs.slide_height = Inches(slide_h)

    # Render each slide from the template
    for slide_tmpl in slides:
        slide_num = slide_tmpl.get("slide_number", "?")
        tmpl_id = slide_tmpl.get("template_id", "unknown")
        obj_count = len(slide_tmpl.get("objects", []))
        print(f"  Slide {slide_num} ({tmpl_id}): {obj_count} objects")
        render_slide(prs, slide_tmpl)

    # Save
    if not output_path:
        OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
        safe_name = topic.replace(" ", "_").replace("/", "-")[:50]
        output_path = str(OUTPUT_DIR / f"{safe_name}.pptx")

    prs.save(output_path)
    print(f"\nSaved: {output_path}")
    print(f"Open in PowerPoint to verify!")
    return output_path


if __name__ == "__main__":
    if len(sys.argv) > 1:
        topic = " ".join(sys.argv[1:])
    else:
        topic = input("Enter presentation topic: ").strip()
        if not topic:
            topic = "Bull Bear Strategy Q1 2025"

    generate_presentation(topic)
