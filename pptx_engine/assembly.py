"""
PPTX Assembly Engine
====================
Core slide-assembly logic for building Potomac PowerPoint decks from
JSON slide plans and branded .pptx template fragments.

Public API:
    assemble_deck(plan, templates_dir, uploads_dir, outputs_dir) -> str
"""

from __future__ import annotations

import copy
import logging
import os
from typing import Any, Dict, Optional

from lxml import etree
from pptx import Presentation
from pptx.util import Emu
from pptx.enum.text import PP_ALIGN

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# 1.  Fill a text placeholder
# ---------------------------------------------------------------------------

def fill_text_placeholder(shape, text: str) -> None:
    """Replace the text of *shape* while preserving the original font
    formatting of the first run in the first paragraph.

    If the shape has no text frame the call is silently ignored.
    """
    if not shape.has_text_frame:
        logger.debug("Shape '%s' has no text frame – skipping text fill.", shape.name)
        return

    tf = shape.text_frame

    # Capture formatting from the first run of the first paragraph
    template_run = None
    if tf.paragraphs and tf.paragraphs[0].runs:
        template_run = tf.paragraphs[0].runs[0]

    # Preserve paragraph-level alignment
    alignment = tf.paragraphs[0].alignment if tf.paragraphs else None

    # Clear all existing paragraphs by removing their XML elements
    p_elements = tf._txBody.findall(
        "{http://schemas.openxmlformats.org/drawingml/2006/main}p"
    )
    for p_el in p_elements:
        tf._txBody.remove(p_el)

    # Split incoming text by newlines so each becomes its own paragraph
    lines = text.split("\n")
    nsmap = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}

    for line in lines:
        p_el = etree.SubElement(tf._txBody, f"{{{nsmap['a']}}}p")
        r_el = etree.SubElement(p_el, f"{{{nsmap['a']}}}r")
        t_el = etree.SubElement(r_el, f"{{{nsmap['a']}}}t")
        t_el.text = line

        # Copy font formatting from template run
        if template_run is not None:
            src_rPr = template_run._r.find(f"{{{nsmap['a']}}}rPr")
            if src_rPr is not None:
                new_rPr = copy.deepcopy(src_rPr)
                r_el.insert(0, new_rPr)

        # Restore alignment
        if alignment is not None:
            pPr = p_el.find(f"{{{nsmap['a']}}}pPr")
            if pPr is None:
                pPr = etree.SubElement(p_el, f"{{{nsmap['a']}}}pPr")
                p_el.insert(0, pPr)
            pPr.set("algn", {
                PP_ALIGN.LEFT: "l",
                PP_ALIGN.CENTER: "ctr",
                PP_ALIGN.RIGHT: "r",
                PP_ALIGN.JUSTIFY: "just",
            }.get(alignment, "l"))

    logger.debug("Filled text shape '%s' with %d line(s).", shape.name, len(lines))


# ---------------------------------------------------------------------------
# 2.  Fill an image placeholder
# ---------------------------------------------------------------------------

def fill_image_placeholder(slide, shape, image_path: str) -> None:
    """Remove *shape* from *slide* and insert an image at the same position
    and size.  Raises ``FileNotFoundError`` if *image_path* doesn't exist.
    """
    if not os.path.isfile(image_path):
        raise FileNotFoundError(f"Image not found: {image_path}")

    # Capture geometry before removing
    left = shape.left
    top = shape.top
    width = shape.width
    height = shape.height

    # Remove the placeholder shape
    sp = shape._element
    sp.getparent().remove(sp)

    # Insert the image in the same bounding box
    slide.shapes.add_picture(image_path, left, top, width, height)

    logger.debug(
        "Replaced shape with image '%s' at (%s, %s, %s, %s).",
        os.path.basename(image_path), left, top, width, height,
    )


# ---------------------------------------------------------------------------
# 3.  Fill an entire slide
# ---------------------------------------------------------------------------

def fill_slide(
    slide,
    content_dict: Dict[str, Any],
    uploaded_images_dict: Optional[Dict[str, str]] = None,
) -> None:
    """Walk every shape on *slide* and fill it from *content_dict*.

    Keys of *content_dict* correspond to shape names in the template.
    Values are either:
      • a **string** → treated as text replacement
      • a **dict** ``{"type": "image", "upload_id": "<id>"}`` → image fill

    *uploaded_images_dict* maps upload_id → absolute path on disk.
    """
    uploaded_images_dict = uploaded_images_dict or {}

    # Build a snapshot list because we mutate the shape tree
    shapes_snapshot = list(slide.shapes)

    for shape in shapes_snapshot:
        name = shape.name
        if name not in content_dict:
            continue

        value = content_dict[name]

        if isinstance(value, str):
            fill_text_placeholder(shape, value)

        elif isinstance(value, dict) and value.get("type") == "image":
            upload_id = value.get("upload_id", "")
            image_path = uploaded_images_dict.get(upload_id, "")
            if not image_path:
                logger.warning(
                    "upload_id '%s' not found in uploaded_images_dict – "
                    "skipping image fill for shape '%s'.",
                    upload_id, name,
                )
                continue
            fill_image_placeholder(slide, shape, image_path)

        else:
            logger.warning(
                "Unsupported content value for shape '%s': %r", name, value
            )


# ---------------------------------------------------------------------------
# 4.  Copy a single slide between presentations
# ---------------------------------------------------------------------------

def copy_slide_into(src_prs: Presentation, slide_idx: int, dest_prs: Presentation):
    """Copy the slide at *slide_idx* in *src_prs* into *dest_prs* and return
    the newly appended slide object.

    Uses deep-copy of the slide XML and relationship parts.
    """
    src_slide = src_prs.slides[slide_idx]

    # Copy the slide layout's slide-layout XML into the destination
    # We create a blank layout from the dest_prs to host the shapes
    slide_layout = dest_prs.slide_layouts[6]  # typically "Blank"
    # Fallback – just pick the first layout if index 6 doesn't exist
    try:
        slide_layout = dest_prs.slide_layouts[6]
    except IndexError:
        slide_layout = dest_prs.slide_layouts[0]

    new_slide = dest_prs.slides.add_slide(slide_layout)

    # ---- copy shapes (deep-copy XML) ----
    # Clear the auto-generated shapes on the blank slide
    for shape_el in list(new_slide.shapes._spTree):
        tag = etree.QName(shape_el.tag).localname
        if tag in ("sp", "pic", "grpSp", "graphicFrame", "cxnSp"):
            new_slide.shapes._spTree.remove(shape_el)

    for shape_el in src_slide.shapes._spTree:
        tag = etree.QName(shape_el.tag).localname
        if tag in ("sp", "pic", "grpSp", "graphicFrame", "cxnSp"):
            new_el = copy.deepcopy(shape_el)
            new_slide.shapes._spTree.append(new_el)

    # ---- copy background ----
    src_bg = src_slide._element.find(
        "{http://schemas.openxmlformats.org/presentationml/2006/main}bg"
    )
    if src_bg is not None:
        # Remove existing background on dest slide
        dest_bg = new_slide._element.find(
            "{http://schemas.openxmlformats.org/presentationml/2006/main}bg"
        )
        if dest_bg is not None:
            new_slide._element.remove(dest_bg)
        new_bg = copy.deepcopy(src_bg)
        # Insert background before the shape tree
        cSld = new_slide._element.find(
            "{http://schemas.openxmlformats.org/presentationml/2006/main}cSld"
        )
        if cSld is not None:
            cSld.insert(0, new_bg)

    # ---- copy images / relationships ----
    for rel in src_slide.part.rels.values():
        if "image" in rel.reltype:
            try:
                new_slide.part.rels.get_or_add(
                    rel.reltype, rel.target_part
                )
            except Exception:
                # Some relationships may not transfer cleanly – log and skip
                logger.debug("Could not copy relationship %s", rel.reltype)

    logger.debug(
        "Copied slide %d from source into destination (now %d slides).",
        slide_idx, len(dest_prs.slides),
    )
    return new_slide


# ---------------------------------------------------------------------------
# 5.  Main assembly function
# ---------------------------------------------------------------------------

def assemble_deck(
    plan: dict,
    templates_dir: str,
    uploads_dir: str,
    outputs_dir: str,
) -> str:
    """Build a complete .pptx deck from a JSON slide plan.

    Parameters
    ----------
    plan : dict
        Expected schema::

            {
                "deck_family": "bull-bear",
                "presentation_title": "Q4 Market Outlook",
                "output_filename": "Q4_Market_Outlook.pptx",
                "slides": [
                    {
                        "template_id": "cover-hero",
                        "content": {
                            "Title 1": "Q4 Market Outlook",
                            "Subtitle 2": "Prepared for Acme Corp",
                            "chart_placeholder": {
                                "type": "image",
                                "upload_id": "abc-123"
                            }
                        },
                        "notes": "Opening remarks …"
                    },
                    ...
                ]
            }

    templates_dir : str
        Root directory containing ``{deck_family}/{template_id}.pptx`` files.
    uploads_dir : str
        Directory where uploaded images are stored.
    outputs_dir : str
        Directory where the finished deck will be saved.

    Returns
    -------
    str
        Absolute path to the saved .pptx file.
    """
    deck_family = plan.get("deck_family", "default")
    output_filename = plan.get("output_filename", "output.pptx")
    slides_plan = plan.get("slides", [])

    if not slides_plan:
        raise ValueError("Slide plan contains no slides.")

    os.makedirs(outputs_dir, exist_ok=True)

    # Build a dict of upload_id -> file path for image fills
    uploaded_images: Dict[str, str] = {}
    if os.path.isdir(uploads_dir):
        for fname in os.listdir(uploads_dir):
            # Files are saved as  {uuid}_{original_name}
            uid = fname.split("_", 1)[0]
            uploaded_images[uid] = os.path.join(uploads_dir, fname)

    dest_prs = Presentation()
    dimensions_set = False

    for idx, slide_spec in enumerate(slides_plan):
        template_id = slide_spec.get("template_id", "blank")
        content = slide_spec.get("content", {})
        notes_text = slide_spec.get("notes")

        # Locate the template file
        template_path = os.path.join(
            templates_dir, deck_family, f"{template_id}.pptx"
        )
        if not os.path.isfile(template_path):
            logger.warning(
                "Template not found: %s – inserting blank slide.", template_path
            )
            layout = dest_prs.slide_layouts[6] if len(dest_prs.slide_layouts) > 6 else dest_prs.slide_layouts[0]
            new_slide = dest_prs.slides.add_slide(layout)
        else:
            src_prs = Presentation(template_path)

            # Set slide dimensions from the first valid template
            if not dimensions_set:
                dest_prs.slide_width = src_prs.slide_width
                dest_prs.slide_height = src_prs.slide_height
                dimensions_set = True
                logger.info(
                    "Slide dimensions set to %s × %s from template '%s'.",
                    src_prs.slide_width, src_prs.slide_height, template_id,
                )

            new_slide = copy_slide_into(src_prs, 0, dest_prs)

        # Fill content
        fill_slide(new_slide, content, uploaded_images)

        # Speaker notes
        if notes_text:
            notes_slide = new_slide.notes_slide
            notes_slide.notes_text_frame.text = notes_text

        logger.info("Assembled slide %d/%d  [%s]", idx + 1, len(slides_plan), template_id)

    # Save
    output_path = os.path.join(outputs_dir, output_filename)
    dest_prs.save(output_path)
    abs_path = os.path.abspath(output_path)
    logger.info("Deck saved → %s  (%d slides)", abs_path, len(slides_plan))
    return abs_path
